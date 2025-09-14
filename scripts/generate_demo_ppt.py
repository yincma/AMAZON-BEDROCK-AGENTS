#!/usr/bin/env python3
"""
AI-PPT-Assistant 演示PPT生成器

功能：
1. 生成包含真实图片的演示PPT
2. 验证图片显示效果
3. 生成质量评估报告
4. 提供用户体验评估

作者: AI-PPT-Assistant Team
日期: 2025-01-14
"""

import os
import sys
import json
import logging
import argparse
import tempfile
import time
from pathlib import Path
from typing import Dict, List, Any, Optional
from io import BytesIO
import base64

# 添加项目路径
project_root = Path(__file__).parent.parent
sys.path.insert(0, str(project_root / "lambdas"))

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.shapes import MSO_SHAPE
    from PIL import Image, ImageDraw, ImageFont
    import boto3
except ImportError as e:
    print(f"警告: 缺少依赖库 {e}")
    print("请安装: pip install python-pptx pillow boto3")
    sys.exit(1)

# 配置日志
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(name)s - %(levelname)s - %(message)s'
)
logger = logging.getLogger(__name__)


class DemoPPTGenerator:
    """演示PPT生成器"""

    def __init__(self, output_dir: str = None):
        """
        初始化生成器

        Args:
            output_dir: 输出目录，默认为临时目录
        """
        self.output_dir = output_dir or tempfile.mkdtemp(prefix="demo_ppt_")
        os.makedirs(self.output_dir, exist_ok=True)

        # 演示数据配置
        self.demo_presentations = {
            "tech_overview": {
                "title": "人工智能技术概览",
                "description": "展示AI技术在各个领域的应用",
                "slides": [
                    {
                        "title": "人工智能发展历程",
                        "content": [
                            "从规则引擎到机器学习",
                            "深度学习的突破性进展",
                            "大语言模型的兴起",
                            "多模态AI的发展趋势"
                        ],
                        "image_prompt": "artificial intelligence development timeline, futuristic technology, digital transformation",
                        "template": "business"
                    },
                    {
                        "title": "机器学习核心算法",
                        "content": [
                            "监督学习：分类与回归",
                            "无监督学习：聚类与降维",
                            "强化学习：决策优化",
                            "深度学习：神经网络架构"
                        ],
                        "image_prompt": "machine learning algorithms visualization, neural networks, data science",
                        "template": "technical"
                    },
                    {
                        "title": "AI在各行业的应用",
                        "content": [
                            "医疗诊断与药物发现",
                            "金融风控与量化交易",
                            "自动驾驶与智能制造",
                            "教育个性化与内容生成"
                        ],
                        "image_prompt": "AI applications across industries, healthcare, finance, automotive, modern technology",
                        "template": "business"
                    }
                ]
            },
            "business_strategy": {
                "title": "数字化转型战略",
                "description": "企业数字化转型的战略规划",
                "slides": [
                    {
                        "title": "数字化转型的必要性",
                        "content": [
                            "市场竞争环境的变化",
                            "客户需求的数字化转变",
                            "技术驱动的商业模式创新",
                            "疫情加速数字化进程"
                        ],
                        "image_prompt": "digital transformation business strategy, corporate innovation, modern office",
                        "template": "business"
                    },
                    {
                        "title": "技术架构设计",
                        "content": [
                            "云原生架构的优势",
                            "微服务与容器化部署",
                            "数据湖与实时分析",
                            "AI/ML平台集成"
                        ],
                        "image_prompt": "cloud architecture diagram, microservices, data analytics platform",
                        "template": "technical"
                    },
                    {
                        "title": "实施路线图",
                        "content": [
                            "第一阶段：基础设施建设",
                            "第二阶段：业务流程数字化",
                            "第三阶段：智能化升级",
                            "第四阶段：生态系统构建"
                        ],
                        "image_prompt": "project roadmap timeline, business transformation phases, strategic planning",
                        "template": "business"
                    }
                ]
            },
            "academic_research": {
                "title": "深度学习研究前沿",
                "description": "学术研究中的深度学习最新进展",
                "slides": [
                    {
                        "title": "Transformer架构创新",
                        "content": [
                            "注意力机制的理论基础",
                            "BERT与GPT系列模型对比",
                            "Vision Transformer在图像处理中的应用",
                            "多模态Transformer的发展"
                        ],
                        "image_prompt": "transformer neural network architecture, attention mechanism visualization, academic research",
                        "template": "academic"
                    },
                    {
                        "title": "生成式AI的突破",
                        "content": [
                            "GAN到Diffusion模型的演进",
                            "文本到图像生成技术",
                            "代码生成与程序合成",
                            "科学发现中的AI应用"
                        ],
                        "image_prompt": "generative AI models, diffusion process, scientific research visualization",
                        "template": "academic"
                    },
                    {
                        "title": "未来研究方向",
                        "content": [
                            "可解释AI与因果推理",
                            "小样本学习与元学习",
                            "神经符号结合",
                            "量子机器学习"
                        ],
                        "image_prompt": "future AI research directions, quantum computing, scientific innovation",
                        "template": "academic"
                    }
                ]
            }
        }

        # 模板样式配置
        self.template_styles = {
            "business": {
                "primary_color": RGBColor(0, 51, 102),  # 深蓝色
                "secondary_color": RGBColor(102, 153, 204),  # 浅蓝色
                "accent_color": RGBColor(255, 102, 0),  # 橙色
                "background_color": RGBColor(255, 255, 255),  # 白色
                "font_name": "Calibri"
            },
            "technical": {
                "primary_color": RGBColor(51, 51, 51),  # 深灰色
                "secondary_color": RGBColor(102, 102, 102),  # 中灰色
                "accent_color": RGBColor(0, 153, 76),  # 绿色
                "background_color": RGBColor(248, 248, 248),  # 浅灰色
                "font_name": "Segoe UI"
            },
            "academic": {
                "primary_color": RGBColor(102, 0, 51),  # 深紫色
                "secondary_color": RGBColor(153, 102, 153),  # 浅紫色
                "accent_color": RGBColor(204, 153, 0),  # 金色
                "background_color": RGBColor(255, 255, 255),  # 白色
                "font_name": "Times New Roman"
            }
        }

    def generate_demo_ppt(self, presentation_type: str = "tech_overview",
                         include_images: bool = True,
                         image_style: str = "placeholder") -> str:
        """
        生成演示PPT

        Args:
            presentation_type: 演示类型 ("tech_overview", "business_strategy", "academic_research")
            include_images: 是否包含图片
            image_style: 图片类型 ("placeholder", "generated", "stock")

        Returns:
            生成的PPT文件路径
        """
        logger.info(f"开始生成演示PPT: {presentation_type}")

        if presentation_type not in self.demo_presentations:
            raise ValueError(f"不支持的演示类型: {presentation_type}")

        presentation_data = self.demo_presentations[presentation_type]

        # 创建PowerPoint演示文稿
        prs = Presentation()

        # 添加标题页
        title_slide = self._create_title_slide(prs, presentation_data)

        # 添加内容页
        for i, slide_data in enumerate(presentation_data["slides"]):
            content_slide = self._create_content_slide(
                prs, slide_data, include_images, image_style
            )

        # 添加总结页
        summary_slide = self._create_summary_slide(prs, presentation_data)

        # 保存PPT
        output_path = os.path.join(
            self.output_dir,
            f"{presentation_type}_demo_{int(time.time())}.pptx"
        )
        prs.save(output_path)

        logger.info(f"PPT生成完成: {output_path}")
        return output_path

    def _create_title_slide(self, prs: Presentation, presentation_data: Dict) -> None:
        """创建标题页"""
        title_layout = prs.slide_layouts[0]  # 标题幻灯片布局
        slide = prs.slides.add_slide(title_layout)

        # 设置标题
        title = slide.shapes.title
        title.text = presentation_data["title"]
        self._apply_title_formatting(title)

        # 设置副标题
        subtitle = slide.placeholders[1]
        subtitle.text = presentation_data["description"]
        self._apply_subtitle_formatting(subtitle)

        # 添加生成信息
        textbox = slide.shapes.add_textbox(
            Inches(7), Inches(6.5), Inches(3), Inches(1)
        )
        text_frame = textbox.text_frame
        p = text_frame.paragraphs[0]
        p.text = f"Generated by AI-PPT-Assistant\n{time.strftime('%Y-%m-%d %H:%M:%S')}"
        p.font.size = Pt(10)
        p.font.color.rgb = RGBColor(128, 128, 128)

    def _create_content_slide(self, prs: Presentation, slide_data: Dict,
                            include_images: bool, image_style: str) -> None:
        """创建内容页"""
        content_layout = prs.slide_layouts[1]  # 标题和内容布局
        slide = prs.slides.add_slide(content_layout)

        # 设置标题
        title = slide.shapes.title
        title.text = slide_data["title"]
        self._apply_content_title_formatting(title)

        # 设置内容
        content_placeholder = slide.placeholders[1]
        text_frame = content_placeholder.text_frame
        text_frame.clear()

        for i, bullet_point in enumerate(slide_data["content"]):
            if i == 0:
                p = text_frame.paragraphs[0]
            else:
                p = text_frame.add_paragraph()

            p.text = bullet_point
            p.level = 0
            self._apply_bullet_formatting(p)

        # 添加图片（如果需要）
        if include_images:
            self._add_image_to_slide(slide, slide_data, image_style)

    def _create_summary_slide(self, prs: Presentation, presentation_data: Dict) -> None:
        """创建总结页"""
        title_layout = prs.slide_layouts[1]  # 标题和内容布局
        slide = prs.slides.add_slide(title_layout)

        # 设置标题
        title = slide.shapes.title
        title.text = "总结与展望"
        self._apply_content_title_formatting(title)

        # 设置总结内容
        content_placeholder = slide.placeholders[1]
        text_frame = content_placeholder.text_frame
        text_frame.clear()

        summary_points = [
            f"✓ {presentation_data['title']}的核心要点",
            "✓ 技术发展趋势分析",
            "✓ 实践应用案例研究",
            "✓ 未来发展方向预测",
            "",
            "感谢观看！"
        ]

        for i, point in enumerate(summary_points):
            if i == 0:
                p = text_frame.paragraphs[0]
            else:
                p = text_frame.add_paragraph()

            p.text = point
            if point == "感谢观看！":
                p.alignment = PP_ALIGN.CENTER
                p.font.size = Pt(24)
                p.font.bold = True
                p.font.color.rgb = RGBColor(0, 51, 102)
            else:
                self._apply_bullet_formatting(p)

    def _add_image_to_slide(self, slide, slide_data: Dict, image_style: str):
        """为幻灯片添加图片"""
        try:
            if image_style == "placeholder":
                image_data = self._create_placeholder_image(slide_data)
            elif image_style == "generated":
                image_data = self._create_ai_generated_image(slide_data)
            else:
                image_data = self._create_stock_image(slide_data)

            if image_data:
                # 添加图片到幻灯片
                img_stream = BytesIO(image_data)
                left = Inches(7.5)
                top = Inches(2)
                width = Inches(2.5)
                height = Inches(3.5)

                slide.shapes.add_picture(img_stream, left, top, width, height)

        except Exception as e:
            logger.warning(f"添加图片失败: {e}")

    def _create_placeholder_image(self, slide_data: Dict) -> bytes:
        """创建占位图"""
        # 创建一个简单的占位图
        img = Image.new('RGB', (800, 600), color=(240, 240, 240))
        draw = ImageDraw.Draw(img)

        # 绘制边框
        draw.rectangle([10, 10, 790, 590], outline=(200, 200, 200), width=3)

        # 添加文本
        try:
            # 尝试使用系统字体
            font = ImageFont.truetype("/System/Library/Fonts/Arial.ttf", 40)
        except:
            font = ImageFont.load_default()

        text = "图片占位符"
        text_bbox = draw.textbbox((0, 0), text, font=font)
        text_width = text_bbox[2] - text_bbox[0]
        text_height = text_bbox[3] - text_bbox[1]

        text_x = (800 - text_width) // 2
        text_y = (600 - text_height) // 2

        draw.text((text_x, text_y), text, fill=(100, 100, 100), font=font)

        # 添加关键词
        keywords = slide_data.get("image_prompt", "").split(",")[0:2]
        if keywords:
            keyword_text = " | ".join(keywords)
            try:
                small_font = ImageFont.truetype("/System/Library/Fonts/Arial.ttf", 20)
            except:
                small_font = ImageFont.load_default()

            keyword_bbox = draw.textbbox((0, 0), keyword_text, font=small_font)
            keyword_width = keyword_bbox[2] - keyword_bbox[0]
            keyword_x = (800 - keyword_width) // 2
            keyword_y = text_y + text_height + 20

            draw.text((keyword_x, keyword_y), keyword_text, fill=(150, 150, 150), font=small_font)

        # 转换为字节
        img_bytes = BytesIO()
        img.save(img_bytes, format='PNG')
        img_bytes.seek(0)
        return img_bytes.getvalue()

    def _create_ai_generated_image(self, slide_data: Dict) -> Optional[bytes]:
        """创建AI生成图片（模拟）"""
        # 这里应该调用真实的AI图片生成服务
        # 现在使用高级占位图模拟
        img = Image.new('RGB', (1024, 768), color=(248, 250, 252))
        draw = ImageDraw.Draw(img)

        # 绘制现代风格的图形
        colors = [
            (59, 130, 246),   # 蓝色
            (16, 185, 129),   # 绿色
            (245, 158, 11),   # 黄色
            (239, 68, 68),    # 红色
        ]

        # 绘制几何形状
        for i, color in enumerate(colors):
            x = 200 + i * 150
            y = 200 + (i % 2) * 100
            draw.ellipse([x, y, x + 100, y + 100], fill=color, outline=color)

        # 添加标题
        try:
            title_font = ImageFont.truetype("/System/Library/Fonts/Arial.ttf", 48)
            subtitle_font = ImageFont.truetype("/System/Library/Fonts/Arial.ttf", 24)
        except:
            title_font = ImageFont.load_default()
            subtitle_font = ImageFont.load_default()

        title_text = "AI Generated"
        draw.text((350, 50), title_text, fill=(31, 41, 55), font=title_font)

        prompt_preview = slide_data.get("image_prompt", "")[:30] + "..."
        draw.text((300, 120), prompt_preview, fill=(107, 114, 128), font=subtitle_font)

        # 转换为字节
        img_bytes = BytesIO()
        img.save(img_bytes, format='PNG')
        img_bytes.seek(0)
        return img_bytes.getvalue()

    def _create_stock_image(self, slide_data: Dict) -> Optional[bytes]:
        """创建股票图片风格（模拟）"""
        # 模拟专业摄影风格的图片
        img = Image.new('RGB', (1200, 800), color=(255, 255, 255))
        draw = ImageDraw.Draw(img)

        # 绘制专业背景
        for y in range(0, 800, 10):
            color_intensity = int(255 - (y / 800) * 50)
            draw.line([(0, y), (1200, y)], fill=(color_intensity, color_intensity, color_intensity))

        # 添加水印
        try:
            watermark_font = ImageFont.truetype("/System/Library/Fonts/Arial.ttf", 36)
        except:
            watermark_font = ImageFont.load_default()

        watermark_text = "Stock Photo Style"
        draw.text((50, 50), watermark_text, fill=(128, 128, 128, 128), font=watermark_font)

        # 转换为字节
        img_bytes = BytesIO()
        img.save(img_bytes, format='PNG')
        img_bytes.seek(0)
        return img_bytes.getvalue()

    def _apply_title_formatting(self, title_shape):
        """应用标题格式"""
        title_shape.text_frame.paragraphs[0].font.size = Pt(44)
        title_shape.text_frame.paragraphs[0].font.bold = True
        title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
        title_shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    def _apply_subtitle_formatting(self, subtitle_shape):
        """应用副标题格式"""
        subtitle_shape.text_frame.paragraphs[0].font.size = Pt(20)
        subtitle_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(102, 102, 102)
        subtitle_shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    def _apply_content_title_formatting(self, title_shape):
        """应用内容页标题格式"""
        title_shape.text_frame.paragraphs[0].font.size = Pt(32)
        title_shape.text_frame.paragraphs[0].font.bold = True
        title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)

    def _apply_bullet_formatting(self, paragraph):
        """应用项目符号格式"""
        paragraph.font.size = Pt(18)
        paragraph.font.color.rgb = RGBColor(51, 51, 51)
        paragraph.space_after = Pt(12)

    def validate_ppt_images(self, ppt_path: str) -> Dict[str, Any]:
        """
        验证PPT中的图片显示效果

        Args:
            ppt_path: PPT文件路径

        Returns:
            验证结果字典
        """
        logger.info(f"开始验证PPT图片: {ppt_path}")

        try:
            prs = Presentation(ppt_path)
            validation_results = {
                'file_path': ppt_path,
                'total_slides': len(prs.slides),
                'slides_with_images': 0,
                'total_images': 0,
                'image_details': [],
                'quality_metrics': {},
                'issues': []
            }

            for slide_idx, slide in enumerate(prs.slides):
                slide_images = []
                for shape in slide.shapes:
                    if hasattr(shape, 'image'):
                        try:
                            image_blob = shape.image.blob
                            img = Image.open(BytesIO(image_blob))

                            image_info = {
                                'slide_number': slide_idx + 1,
                                'format': img.format,
                                'size': img.size,
                                'mode': img.mode,
                                'file_size': len(image_blob),
                                'position': (shape.left, shape.top),
                                'dimensions': (shape.width, shape.height)
                            }

                            slide_images.append(image_info)
                            validation_results['total_images'] += 1

                        except Exception as e:
                            validation_results['issues'].append({
                                'slide': slide_idx + 1,
                                'issue': f'图片读取失败: {str(e)}'
                            })

                if slide_images:
                    validation_results['slides_with_images'] += 1
                    validation_results['image_details'].extend(slide_images)

            # 计算质量指标
            validation_results['quality_metrics'] = self._calculate_image_quality_metrics(
                validation_results['image_details']
            )

            # 生成验证报告
            self._generate_image_validation_report(validation_results)

            logger.info("PPT图片验证完成")
            return validation_results

        except Exception as e:
            logger.error(f"PPT验证失败: {e}")
            return {
                'error': str(e),
                'file_path': ppt_path
            }

    def _calculate_image_quality_metrics(self, image_details: List[Dict]) -> Dict:
        """计算图片质量指标"""
        if not image_details:
            return {}

        total_images = len(image_details)
        total_file_size = sum(img['file_size'] for img in image_details)
        avg_file_size = total_file_size / total_images

        resolutions = [img['size'] for img in image_details]
        avg_width = sum(r[0] for r in resolutions) / total_images
        avg_height = sum(r[1] for r in resolutions) / total_images

        format_distribution = {}
        for img in image_details:
            fmt = img['format']
            format_distribution[fmt] = format_distribution.get(fmt, 0) + 1

        return {
            'total_images': total_images,
            'average_file_size': avg_file_size,
            'average_resolution': (int(avg_width), int(avg_height)),
            'format_distribution': format_distribution,
            'size_range': {
                'min_width': min(r[0] for r in resolutions),
                'max_width': max(r[0] for r in resolutions),
                'min_height': min(r[1] for r in resolutions),
                'max_height': max(r[1] for r in resolutions)
            }
        }

    def _generate_image_validation_report(self, validation_results: Dict):
        """生成图片验证报告"""
        report_path = os.path.join(self.output_dir, "image_validation_report.json")

        with open(report_path, 'w', encoding='utf-8') as f:
            json.dump(validation_results, f, indent=2, ensure_ascii=False)

        # 生成Markdown报告
        md_report_path = os.path.join(self.output_dir, "image_validation_report.md")

        with open(md_report_path, 'w', encoding='utf-8') as f:
            f.write("# PPT图片验证报告\n\n")
            f.write(f"文件: {validation_results['file_path']}\n")
            f.write(f"生成时间: {time.strftime('%Y-%m-%d %H:%M:%S')}\n\n")

            f.write("## 基本信息\n")
            f.write(f"- 总幻灯片数: {validation_results['total_slides']}\n")
            f.write(f"- 包含图片的幻灯片数: {validation_results['slides_with_images']}\n")
            f.write(f"- 总图片数: {validation_results['total_images']}\n\n")

            if validation_results.get('quality_metrics'):
                metrics = validation_results['quality_metrics']
                f.write("## 质量指标\n")
                f.write(f"- 平均文件大小: {metrics['average_file_size']:.0f} 字节\n")
                f.write(f"- 平均分辨率: {metrics['average_resolution'][0]}×{metrics['average_resolution'][1]}\n")
                f.write(f"- 格式分布: {metrics['format_distribution']}\n\n")

            if validation_results.get('issues'):
                f.write("## 发现的问题\n")
                for issue in validation_results['issues']:
                    f.write(f"- 幻灯片 {issue['slide']}: {issue['issue']}\n")
                f.write("\n")

            f.write("## 改进建议\n")
            f.write("1. 确保图片分辨率适合演示使用\n")
            f.write("2. 优化图片文件大小以提高加载速度\n")
            f.write("3. 保持图片风格的一致性\n")
            f.write("4. 验证图片与内容的相关性\n")

        logger.info(f"验证报告已生成: {md_report_path}")

    def run_comprehensive_test(self) -> Dict[str, Any]:
        """运行综合测试"""
        logger.info("开始运行综合PPT图片显示效果测试...")

        test_results = {
            'timestamp': time.strftime('%Y-%m-%d %H:%M:%S'),
            'presentations': {},
            'summary': {}
        }

        # 测试所有演示类型
        for ppt_type in self.demo_presentations.keys():
            logger.info(f"测试演示类型: {ppt_type}")

            try:
                # 生成PPT
                ppt_path = self.generate_demo_ppt(
                    presentation_type=ppt_type,
                    include_images=True,
                    image_style="placeholder"
                )

                # 验证图片
                validation_results = self.validate_ppt_images(ppt_path)

                test_results['presentations'][ppt_type] = {
                    'ppt_path': ppt_path,
                    'validation': validation_results,
                    'status': 'success'
                }

            except Exception as e:
                logger.error(f"测试 {ppt_type} 失败: {e}")
                test_results['presentations'][ppt_type] = {
                    'status': 'error',
                    'error': str(e)
                }

        # 生成综合摘要
        test_results['summary'] = self._generate_test_summary(test_results['presentations'])

        # 保存综合报告
        self._save_comprehensive_report(test_results)

        return test_results

    def _generate_test_summary(self, presentations: Dict) -> Dict:
        """生成测试摘要"""
        total_tests = len(presentations)
        successful_tests = len([p for p in presentations.values() if p.get('status') == 'success'])
        total_images = sum(
            p.get('validation', {}).get('total_images', 0)
            for p in presentations.values()
            if p.get('status') == 'success'
        )

        return {
            'total_presentations_tested': total_tests,
            'successful_tests': successful_tests,
            'success_rate': successful_tests / total_tests if total_tests > 0 else 0,
            'total_images_validated': total_images,
            'test_status': 'PASSED' if successful_tests == total_tests else 'PARTIAL' if successful_tests > 0 else 'FAILED'
        }

    def _save_comprehensive_report(self, test_results: Dict):
        """保存综合报告"""
        report_path = os.path.join(self.output_dir, "comprehensive_test_report.json")

        with open(report_path, 'w', encoding='utf-8') as f:
            json.dump(test_results, f, indent=2, ensure_ascii=False)

        logger.info(f"综合测试报告已保存: {report_path}")


def main():
    """主函数"""
    parser = argparse.ArgumentParser(description="AI-PPT-Assistant 演示PPT生成器")
    parser.add_argument(
        '--type',
        choices=['tech_overview', 'business_strategy', 'academic_research', 'all'],
        default='tech_overview',
        help='生成的演示类型'
    )
    parser.add_argument(
        '--images',
        choices=['placeholder', 'generated', 'stock'],
        default='placeholder',
        help='图片类型'
    )
    parser.add_argument(
        '--output-dir',
        type=str,
        help='输出目录'
    )
    parser.add_argument(
        '--validate-only',
        action='store_true',
        help='仅验证现有PPT文件'
    )
    parser.add_argument(
        '--ppt-file',
        type=str,
        help='要验证的PPT文件路径'
    )
    parser.add_argument(
        '--comprehensive-test',
        action='store_true',
        help='运行综合测试'
    )

    args = parser.parse_args()

    try:
        generator = DemoPPTGenerator(output_dir=args.output_dir)

        if args.validate_only and args.ppt_file:
            # 仅验证指定PPT文件
            logger.info(f"验证PPT文件: {args.ppt_file}")
            results = generator.validate_ppt_images(args.ppt_file)
            print(json.dumps(results, indent=2, ensure_ascii=False))

        elif args.comprehensive_test:
            # 运行综合测试
            results = generator.run_comprehensive_test()
            print(f"综合测试完成!")
            print(f"成功率: {results['summary']['success_rate']:.1%}")
            print(f"状态: {results['summary']['test_status']}")
            print(f"输出目录: {generator.output_dir}")

        elif args.type == 'all':
            # 生成所有类型的演示
            for ppt_type in generator.demo_presentations.keys():
                ppt_path = generator.generate_demo_ppt(
                    presentation_type=ppt_type,
                    include_images=True,
                    image_style=args.images
                )
                print(f"生成完成: {ppt_path}")

        else:
            # 生成指定类型的演示
            ppt_path = generator.generate_demo_ppt(
                presentation_type=args.type,
                include_images=True,
                image_style=args.images
            )
            print(f"演示PPT生成完成: {ppt_path}")

            # 自动验证
            validation_results = generator.validate_ppt_images(ppt_path)
            print(f"图片验证完成: {validation_results['total_images']} 张图片")

    except Exception as e:
        logger.error(f"执行失败: {e}")
        sys.exit(1)


if __name__ == "__main__":
    main()