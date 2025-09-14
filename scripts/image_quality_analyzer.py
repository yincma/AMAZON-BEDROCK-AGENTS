#!/usr/bin/env python3
"""
AI-PPT-Assistant 图片质量分析器

功能：
1. 分析PPT中图片的质量指标
2. 评估图片与内容的相关性
3. 检查图片的技术规范
4. 生成详细的质量评估报告
5. 提供优化建议

作者: AI-PPT-Assistant Team
日期: 2025-01-14
"""

import os
import sys
import json
import logging
import argparse
import tempfile
import time
import hashlib
from pathlib import Path
from typing import Dict, List, Any, Optional, Tuple
from io import BytesIO
import math
import statistics

# 添加项目路径
project_root = Path(__file__).parent.parent
sys.path.insert(0, str(project_root / "lambdas"))

try:
    from PIL import Image, ImageStat, ImageFilter
    from PIL.ExifTags import TAGS
    from pptx import Presentation
    import numpy as np
except ImportError as e:
    print(f"警告: 缺少依赖库 {e}")
    print("请安装: pip install python-pptx pillow numpy")
    sys.exit(1)

# 配置日志
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(name)s - %(levelname)s - %(message)s'
)
logger = logging.getLogger(__name__)


class ImageQualityAnalyzer:
    """图片质量分析器"""

    def __init__(self, output_dir: str = None):
        """
        初始化分析器

        Args:
            output_dir: 输出目录，默认为临时目录
        """
        self.output_dir = output_dir or tempfile.mkdtemp(prefix="image_quality_")
        os.makedirs(self.output_dir, exist_ok=True)

        # 质量标准配置
        self.quality_standards = {
            'presentation': {
                'min_width': 800,
                'min_height': 600,
                'max_width': 3840,
                'max_height': 2160,
                'min_dpi': 96,
                'recommended_dpi': 150,
                'max_file_size': 5 * 1024 * 1024,  # 5MB
                'recommended_file_size': 2 * 1024 * 1024,  # 2MB
                'aspect_ratios': [(16, 9), (4, 3), (3, 2), (1, 1)],
                'supported_formats': ['PNG', 'JPEG', 'JPG', 'GIF', 'BMP'],
                'preferred_formats': ['PNG', 'JPEG']
            },
            'quality_metrics': {
                'min_sharpness': 0.3,
                'min_contrast': 0.2,
                'min_brightness': 0.1,
                'max_brightness': 0.9,
                'min_saturation': 0.1,
                'max_noise_level': 0.3
            }
        }

        # 图片类型检测模式
        self.content_keywords = {
            'business': ['office', 'meeting', 'corporate', 'professional', 'suit', 'handshake'],
            'technology': ['computer', 'digital', 'network', 'data', 'ai', 'robot', 'circuit'],
            'academic': ['book', 'research', 'study', 'library', 'education', 'science'],
            'medical': ['hospital', 'doctor', 'medical', 'health', 'care', 'treatment'],
            'nature': ['landscape', 'tree', 'mountain', 'ocean', 'sky', 'animal'],
            'abstract': ['pattern', 'geometric', 'artistic', 'creative', 'design']
        }

    def analyze_ppt_images(self, ppt_path: str) -> Dict[str, Any]:
        """
        分析PPT中所有图片的质量

        Args:
            ppt_path: PPT文件路径

        Returns:
            分析结果字典
        """
        logger.info(f"开始分析PPT图片质量: {ppt_path}")

        try:
            prs = Presentation(ppt_path)
            analysis_results = {
                'file_info': {
                    'path': ppt_path,
                    'total_slides': len(prs.slides),
                    'analysis_time': time.strftime('%Y-%m-%d %H:%M:%S')
                },
                'image_analysis': [],
                'quality_summary': {},
                'recommendations': [],
                'technical_issues': []
            }

            image_count = 0

            for slide_idx, slide in enumerate(prs.slides, 1):
                for shape_idx, shape in enumerate(slide.shapes):
                    if hasattr(shape, 'image'):
                        try:
                            image_blob = shape.image.blob
                            img = Image.open(BytesIO(image_blob))

                            # 分析单个图片
                            image_analysis = self._analyze_single_image(
                                img, image_blob, slide_idx, shape_idx, shape
                            )

                            analysis_results['image_analysis'].append(image_analysis)
                            image_count += 1

                        except Exception as e:
                            logger.warning(f"分析图片失败 (幻灯片{slide_idx}, 形状{shape_idx}): {e}")
                            analysis_results['technical_issues'].append({
                                'slide': slide_idx,
                                'shape': shape_idx,
                                'issue': f'图片读取失败: {str(e)}'
                            })

            # 生成质量摘要
            analysis_results['quality_summary'] = self._generate_quality_summary(
                analysis_results['image_analysis']
            )

            # 生成优化建议
            analysis_results['recommendations'] = self._generate_recommendations(
                analysis_results['image_analysis'],
                analysis_results['quality_summary']
            )

            # 保存分析报告
            self._save_analysis_report(analysis_results)

            logger.info(f"图片质量分析完成，共分析 {image_count} 张图片")
            return analysis_results

        except Exception as e:
            logger.error(f"PPT分析失败: {e}")
            return {
                'error': str(e),
                'file_path': ppt_path
            }

    def _analyze_single_image(self, img: Image.Image, image_blob: bytes,
                            slide_idx: int, shape_idx: int, shape) -> Dict[str, Any]:
        """分析单个图片"""
        analysis = {
            'location': {
                'slide': slide_idx,
                'shape': shape_idx,
                'position': (int(shape.left), int(shape.top)),
                'size_in_ppt': (int(shape.width), int(shape.height))
            },
            'basic_info': self._get_basic_image_info(img, image_blob),
            'quality_metrics': self._calculate_quality_metrics(img),
            'technical_analysis': self._perform_technical_analysis(img),
            'content_analysis': self._analyze_image_content(img),
            'compliance_check': self._check_compliance(img, image_blob)
        }

        return analysis

    def _get_basic_image_info(self, img: Image.Image, image_blob: bytes) -> Dict[str, Any]:
        """获取图片基本信息"""
        return {
            'format': img.format,
            'mode': img.mode,
            'size': img.size,
            'width': img.width,
            'height': img.height,
            'aspect_ratio': round(img.width / img.height, 2),
            'file_size': len(image_blob),
            'file_size_mb': round(len(image_blob) / (1024 * 1024), 2),
            'has_transparency': img.mode in ('RGBA', 'LA') or 'transparency' in img.info,
            'color_channels': len(img.getbands()) if hasattr(img, 'getbands') else 0
        }

    def _calculate_quality_metrics(self, img: Image.Image) -> Dict[str, Any]:
        """计算图片质量指标"""
        try:
            # 转换为RGB进行分析（如果需要）
            if img.mode != 'RGB':
                analysis_img = img.convert('RGB')
            else:
                analysis_img = img

            # 计算统计信息
            stat = ImageStat.Stat(analysis_img)

            # 亮度分析
            brightness = sum(stat.mean) / (len(stat.mean) * 255)

            # 对比度分析
            contrast = sum(stat.stddev) / (len(stat.stddev) * 255)

            # 锐度分析（使用拉普拉斯算子）
            sharpness = self._calculate_sharpness(analysis_img)

            # 饱和度分析
            hsv_img = analysis_img.convert('HSV')
            hsv_stat = ImageStat.Stat(hsv_img)
            saturation = hsv_stat.mean[1] / 255 if len(hsv_stat.mean) > 1 else 0

            # 噪声水平估计
            noise_level = self._estimate_noise_level(analysis_img)

            # 色彩丰富度
            color_richness = self._calculate_color_richness(analysis_img)

            return {
                'brightness': round(brightness, 3),
                'contrast': round(contrast, 3),
                'sharpness': round(sharpness, 3),
                'saturation': round(saturation, 3),
                'noise_level': round(noise_level, 3),
                'color_richness': round(color_richness, 3),
                'overall_quality_score': self._calculate_overall_quality_score(
                    brightness, contrast, sharpness, saturation, noise_level
                )
            }

        except Exception as e:
            logger.warning(f"质量指标计算失败: {e}")
            return {
                'error': str(e)
            }

    def _calculate_sharpness(self, img: Image.Image) -> float:
        """计算图片锐度"""
        try:
            # 转换为灰度图
            gray = img.convert('L')

            # 应用拉普拉斯算子
            laplacian = gray.filter(ImageFilter.Kernel((3, 3), [-1, -1, -1, -1, 8, -1, -1, -1, -1], 1, 0))

            # 计算方差作为锐度指标
            stat = ImageStat.Stat(laplacian)
            return stat.stddev[0] / 255

        except Exception:
            return 0.0

    def _estimate_noise_level(self, img: Image.Image) -> float:
        """估计图片噪声水平"""
        try:
            # 使用高通滤波器检测噪声
            gray = img.convert('L')

            # 应用高通滤波
            high_pass = gray.filter(ImageFilter.Kernel((3, 3), [-1, -1, -1, -1, 8, -1, -1, -1, -1], 1, 0))

            # 计算高频成分的强度
            stat = ImageStat.Stat(high_pass)
            noise_estimate = stat.mean[0] / 255

            return min(noise_estimate, 1.0)

        except Exception:
            return 0.0

    def _calculate_color_richness(self, img: Image.Image) -> float:
        """计算色彩丰富度"""
        try:
            # 统计不同颜色的数量
            colors = img.getcolors(maxcolors=256*256*256)
            if colors is None:
                # 如果颜色太多，使用近似方法
                return 1.0

            unique_colors = len(colors)
            total_pixels = img.width * img.height

            # 计算色彩丰富度（归一化）
            richness = min(unique_colors / (total_pixels * 0.1), 1.0)
            return richness

        except Exception:
            return 0.5

    def _calculate_overall_quality_score(self, brightness: float, contrast: float,
                                       sharpness: float, saturation: float,
                                       noise_level: float) -> float:
        """计算综合质量分数"""
        try:
            # 亮度评分 (0.2-0.8为最佳)
            brightness_score = 1.0 - abs(brightness - 0.5) * 2

            # 对比度评分 (越高越好，但不超过1)
            contrast_score = min(contrast * 2, 1.0)

            # 锐度评分 (越高越好，但不超过1)
            sharpness_score = min(sharpness * 2, 1.0)

            # 饱和度评分 (0.3-0.8为最佳)
            saturation_score = 1.0 - abs(saturation - 0.55) * 2

            # 噪声评分 (越低越好)
            noise_score = 1.0 - noise_level

            # 加权平均
            weights = {
                'brightness': 0.15,
                'contrast': 0.25,
                'sharpness': 0.30,
                'saturation': 0.15,
                'noise': 0.15
            }

            overall_score = (
                brightness_score * weights['brightness'] +
                contrast_score * weights['contrast'] +
                sharpness_score * weights['sharpness'] +
                saturation_score * weights['saturation'] +
                noise_score * weights['noise']
            )

            return round(max(0, min(1, overall_score)), 3)

        except Exception:
            return 0.5

    def _perform_technical_analysis(self, img: Image.Image) -> Dict[str, Any]:
        """执行技术分析"""
        analysis = {
            'dpi': self._estimate_dpi(img),
            'compression_artifacts': self._detect_compression_artifacts(img),
            'bit_depth': self._estimate_bit_depth(img),
            'color_space': img.mode,
            'has_exif': bool(getattr(img, '_getexif', None) and img._getexif()),
            'potential_upscaling': self._detect_upscaling(img)
        }

        return analysis

    def _estimate_dpi(self, img: Image.Image) -> Optional[Tuple[float, float]]:
        """估计图片DPI"""
        try:
            if hasattr(img, 'info') and 'dpi' in img.info:
                return img.info['dpi']
            else:
                # 默认DPI (常见的屏幕显示DPI)
                return (96.0, 96.0)
        except Exception:
            return None

    def _detect_compression_artifacts(self, img: Image.Image) -> Dict[str, Any]:
        """检测压缩伪影"""
        try:
            if img.format == 'JPEG':
                # JPEG压缩伪影检测
                # 检查8x8块效应
                gray = img.convert('L')
                # 简化的块效应检测
                block_artifacts = self._detect_blocking_artifacts(gray)

                return {
                    'type': 'JPEG',
                    'blocking_artifacts': block_artifacts,
                    'severity': 'low' if block_artifacts < 0.1 else 'medium' if block_artifacts < 0.3 else 'high'
                }
            else:
                return {
                    'type': img.format,
                    'artifacts_detected': False
                }

        except Exception:
            return {'error': 'compression analysis failed'}

    def _detect_blocking_artifacts(self, gray_img: Image.Image) -> float:
        """检测JPEG块效应"""
        try:
            # 简化的块效应检测算法
            width, height = gray_img.size
            if width < 16 or height < 16:
                return 0.0

            # 取样检测（每8像素检查一次边界）
            total_edge_strength = 0
            edge_count = 0

            for y in range(8, height-8, 8):
                for x in range(0, width-1):
                    # 检查垂直边界
                    pixel1 = gray_img.getpixel((x, y-1))
                    pixel2 = gray_img.getpixel((x, y))
                    edge_strength = abs(pixel1 - pixel2)
                    total_edge_strength += edge_strength
                    edge_count += 1

            if edge_count > 0:
                avg_edge_strength = total_edge_strength / edge_count
                return min(avg_edge_strength / 255, 1.0)

            return 0.0

        except Exception:
            return 0.0

    def _estimate_bit_depth(self, img: Image.Image) -> int:
        """估计位深度"""
        try:
            if img.mode == 'L':
                return 8  # 灰度图
            elif img.mode in ['RGB', 'RGBA']:
                return 24 if img.mode == 'RGB' else 32
            elif img.mode == 'P':
                return 8  # 调色板模式
            else:
                return 8  # 默认

        except Exception:
            return 8

    def _detect_upscaling(self, img: Image.Image) -> Dict[str, Any]:
        """检测图片是否被放大"""
        try:
            # 通过分析高频细节检测放大
            gray = img.convert('L')

            # 应用高通滤波器
            high_pass = gray.filter(ImageFilter.FIND_EDGES)
            stat = ImageStat.Stat(high_pass)

            edge_density = stat.mean[0] / 255
            edge_variance = stat.stddev[0] / 255

            # 低边缘密度和方差可能表示放大
            upscaling_likelihood = 0.0
            if edge_density < 0.1 and edge_variance < 0.2:
                upscaling_likelihood = 0.8
            elif edge_density < 0.2 and edge_variance < 0.4:
                upscaling_likelihood = 0.5

            return {
                'likely_upscaled': upscaling_likelihood > 0.6,
                'confidence': upscaling_likelihood,
                'edge_density': round(edge_density, 3),
                'edge_variance': round(edge_variance, 3)
            }

        except Exception:
            return {'error': 'upscaling detection failed'}

    def _analyze_image_content(self, img: Image.Image) -> Dict[str, Any]:
        """分析图片内容类型"""
        try:
            # 基于颜色分布分析内容类型
            rgb_img = img.convert('RGB')
            stat = ImageStat.Stat(rgb_img)

            # 计算主要颜色
            dominant_colors = self._extract_dominant_colors(rgb_img)

            # 分析色彩分布特征
            color_features = {
                'mean_rgb': [round(m, 1) for m in stat.mean],
                'std_rgb': [round(s, 1) for s in stat.stddev],
                'dominant_colors': dominant_colors
            }

            # 估计内容类型
            content_type = self._classify_content_type(color_features)

            return {
                'estimated_type': content_type,
                'color_features': color_features,
                'is_grayscale': self._is_grayscale(rgb_img),
                'has_people': self._detect_skin_tones(rgb_img),
                'is_screenshot': self._detect_screenshot_characteristics(rgb_img)
            }

        except Exception as e:
            return {'error': f'content analysis failed: {str(e)}'}

    def _extract_dominant_colors(self, img: Image.Image, num_colors: int = 5) -> List[Tuple[int, int, int]]:
        """提取主要颜色"""
        try:
            # 缩小图片以提高性能
            img_small = img.resize((100, 100))

            # 获取颜色分布
            colors = img_small.getcolors(maxcolors=10000)
            if colors is None:
                return []

            # 按频率排序
            colors.sort(key=lambda x: x[0], reverse=True)

            # 返回前N个颜色
            return [color[1] for color in colors[:num_colors]]

        except Exception:
            return []

    def _classify_content_type(self, color_features: Dict) -> str:
        """基于颜色特征分类内容类型"""
        try:
            mean_rgb = color_features['mean_rgb']
            std_rgb = color_features['std_rgb']

            # 计算平均亮度
            avg_brightness = sum(mean_rgb) / 3

            # 计算色彩变化
            color_variance = sum(std_rgb) / 3

            # 基于简单规则进行分类
            if avg_brightness > 200 and color_variance < 30:
                return 'document'  # 文档类（高亮度，低变化）
            elif avg_brightness < 100:
                return 'dark_theme'  # 深色主题
            elif color_variance > 80:
                return 'natural_photo'  # 自然照片（高变化）
            elif 100 < avg_brightness < 200 and 30 < color_variance < 80:
                return 'illustration'  # 插图
            else:
                return 'mixed'  # 混合类型

        except Exception:
            return 'unknown'

    def _is_grayscale(self, img: Image.Image) -> bool:
        """检测是否为灰度图"""
        try:
            stat = ImageStat.Stat(img)
            # 检查RGB通道的标准差
            r_std, g_std, b_std = stat.stddev[:3] if len(stat.stddev) >= 3 else (0, 0, 0)

            # 如果各通道差异很小，可能是灰度图
            return abs(r_std - g_std) < 5 and abs(g_std - b_std) < 5 and abs(r_std - b_std) < 5

        except Exception:
            return False

    def _detect_skin_tones(self, img: Image.Image) -> bool:
        """检测肤色（可能包含人物）"""
        try:
            # 肤色范围（RGB）
            skin_ranges = [
                ((95, 40, 20), (255, 219, 172)),  # 浅肤色
                ((80, 25, 10), (200, 150, 100)),  # 中等肤色
                ((30, 15, 5), (120, 80, 60))      # 深肤色
            ]

            pixels = list(img.getdata())
            total_pixels = len(pixels)
            skin_pixels = 0

            for pixel in pixels[:min(1000, total_pixels)]:  # 采样检查
                r, g, b = pixel[:3]
                for (min_r, min_g, min_b), (max_r, max_g, max_b) in skin_ranges:
                    if min_r <= r <= max_r and min_g <= g <= max_g and min_b <= b <= max_b:
                        skin_pixels += 1
                        break

            # 如果肤色像素超过5%，可能包含人物
            return (skin_pixels / min(1000, total_pixels)) > 0.05

        except Exception:
            return False

    def _detect_screenshot_characteristics(self, img: Image.Image) -> bool:
        """检测截图特征"""
        try:
            # 截图通常有锐利的边缘和平坦的区域
            gray = img.convert('L')

            # 检测锐利边缘
            edges = gray.filter(ImageFilter.FIND_EDGES)
            edge_stat = ImageStat.Stat(edges)

            # 检测平坦区域
            smooth = gray.filter(ImageFilter.SMOOTH)
            variance = ImageStat.Stat(smooth).stddev[0]

            # 截图特征：边缘锐利但整体方差较小
            edge_strength = edge_stat.mean[0]
            is_screenshot = edge_strength > 20 and variance < 40

            return is_screenshot

        except Exception:
            return False

    def _check_compliance(self, img: Image.Image, image_blob: bytes) -> Dict[str, Any]:
        """检查是否符合标准"""
        standards = self.quality_standards['presentation']
        compliance = {
            'overall_compliant': True,
            'issues': [],
            'warnings': [],
            'recommendations': []
        }

        # 检查尺寸
        if img.width < standards['min_width'] or img.height < standards['min_height']:
            compliance['overall_compliant'] = False
            compliance['issues'].append(
                f"分辨率过低 ({img.width}x{img.height})，"
                f"建议最小分辨率 {standards['min_width']}x{standards['min_height']}"
            )

        if img.width > standards['max_width'] or img.height > standards['max_height']:
            compliance['warnings'].append(
                f"分辨率过高 ({img.width}x{img.height})，"
                f"可能影响加载速度"
            )

        # 检查文件大小
        file_size = len(image_blob)
        if file_size > standards['max_file_size']:
            compliance['overall_compliant'] = False
            compliance['issues'].append(
                f"文件过大 ({file_size/1024/1024:.1f}MB)，"
                f"超过最大限制 {standards['max_file_size']/1024/1024:.1f}MB"
            )
        elif file_size > standards['recommended_file_size']:
            compliance['warnings'].append(
                f"文件较大 ({file_size/1024/1024:.1f}MB)，"
                f"建议压缩至 {standards['recommended_file_size']/1024/1024:.1f}MB 以下"
            )

        # 检查格式
        if img.format not in standards['supported_formats']:
            compliance['overall_compliant'] = False
            compliance['issues'].append(
                f"不支持的格式 ({img.format})，"
                f"支持的格式: {', '.join(standards['supported_formats'])}"
            )
        elif img.format not in standards['preferred_formats']:
            compliance['recommendations'].append(
                f"建议使用更优格式，当前: {img.format}，"
                f"推荐: {', '.join(standards['preferred_formats'])}"
            )

        return compliance

    def _generate_quality_summary(self, image_analyses: List[Dict]) -> Dict[str, Any]:
        """生成质量摘要"""
        if not image_analyses:
            return {}

        # 提取质量指标
        quality_scores = []
        brightness_values = []
        contrast_values = []
        sharpness_values = []
        file_sizes = []
        resolutions = []

        for analysis in image_analyses:
            quality_metrics = analysis.get('quality_metrics', {})
            if 'overall_quality_score' in quality_metrics:
                quality_scores.append(quality_metrics['overall_quality_score'])
                brightness_values.append(quality_metrics.get('brightness', 0))
                contrast_values.append(quality_metrics.get('contrast', 0))
                sharpness_values.append(quality_metrics.get('sharpness', 0))

            basic_info = analysis.get('basic_info', {})
            if 'file_size' in basic_info:
                file_sizes.append(basic_info['file_size'])

            if 'size' in basic_info:
                resolutions.append(basic_info['size'])

        # 计算统计信息
        summary = {
            'total_images': len(image_analyses),
            'quality_statistics': {},
            'technical_summary': {},
            'compliance_summary': {}
        }

        if quality_scores:
            summary['quality_statistics'] = {
                'average_quality_score': round(statistics.mean(quality_scores), 3),
                'min_quality_score': round(min(quality_scores), 3),
                'max_quality_score': round(max(quality_scores), 3),
                'quality_distribution': self._categorize_quality_scores(quality_scores)
            }

            summary['technical_summary'] = {
                'average_brightness': round(statistics.mean(brightness_values), 3),
                'average_contrast': round(statistics.mean(contrast_values), 3),
                'average_sharpness': round(statistics.mean(sharpness_values), 3)
            }

        if file_sizes:
            summary['technical_summary'].update({
                'total_file_size': sum(file_sizes),
                'average_file_size': round(statistics.mean(file_sizes), 0),
                'largest_file': max(file_sizes),
                'smallest_file': min(file_sizes)
            })

        # 合规性统计
        compliant_images = sum(1 for analysis in image_analyses
                             if analysis.get('compliance_check', {}).get('overall_compliant', False))

        summary['compliance_summary'] = {
            'compliant_images': compliant_images,
            'compliance_rate': round(compliant_images / len(image_analyses), 3) if image_analyses else 0,
            'total_issues': sum(len(analysis.get('compliance_check', {}).get('issues', []))
                              for analysis in image_analyses),
            'total_warnings': sum(len(analysis.get('compliance_check', {}).get('warnings', []))
                                for analysis in image_analyses)
        }

        return summary

    def _categorize_quality_scores(self, scores: List[float]) -> Dict[str, int]:
        """分类质量分数"""
        categories = {'excellent': 0, 'good': 0, 'fair': 0, 'poor': 0}

        for score in scores:
            if score >= 0.8:
                categories['excellent'] += 1
            elif score >= 0.6:
                categories['good'] += 1
            elif score >= 0.4:
                categories['fair'] += 1
            else:
                categories['poor'] += 1

        return categories

    def _generate_recommendations(self, image_analyses: List[Dict],
                                quality_summary: Dict) -> List[str]:
        """生成优化建议"""
        recommendations = []

        # 基于质量摘要的建议
        quality_stats = quality_summary.get('quality_statistics', {})
        avg_quality = quality_stats.get('average_quality_score', 0)

        if avg_quality < 0.5:
            recommendations.append("整体图片质量较低，建议使用更高质量的原始图片")

        quality_dist = quality_stats.get('quality_distribution', {})
        if quality_dist.get('poor', 0) > 0:
            recommendations.append(f"发现 {quality_dist['poor']} 张质量较差的图片，建议重新生成或替换")

        # 基于技术指标的建议
        tech_summary = quality_summary.get('technical_summary', {})

        avg_brightness = tech_summary.get('average_brightness', 0.5)
        if avg_brightness < 0.3:
            recommendations.append("图片整体偏暗，建议增加亮度以提高可读性")
        elif avg_brightness > 0.8:
            recommendations.append("图片整体过亮，建议降低亮度以避免刺眼")

        avg_contrast = tech_summary.get('average_contrast', 0.5)
        if avg_contrast < 0.3:
            recommendations.append("图片对比度不足，建议增强对比度以提高清晰度")

        avg_sharpness = tech_summary.get('average_sharpness', 0.5)
        if avg_sharpness < 0.3:
            recommendations.append("图片锐度不足，建议使用更清晰的原始图片")

        # 文件大小建议
        avg_file_size = tech_summary.get('average_file_size', 0)
        if avg_file_size > 2 * 1024 * 1024:  # 2MB
            recommendations.append("图片文件较大，建议压缩以提高加载速度")

        # 合规性建议
        compliance_summary = quality_summary.get('compliance_summary', {})
        compliance_rate = compliance_summary.get('compliance_rate', 1)

        if compliance_rate < 1.0:
            recommendations.append(f"有 {int((1-compliance_rate)*100)}% 的图片不符合标准，请检查具体问题")

        total_issues = compliance_summary.get('total_issues', 0)
        if total_issues > 0:
            recommendations.append(f"发现 {total_issues} 个合规性问题，需要立即修复")

        total_warnings = compliance_summary.get('total_warnings', 0)
        if total_warnings > 0:
            recommendations.append(f"有 {total_warnings} 个警告项，建议优化以获得更好效果")

        # 通用建议
        if not recommendations:
            recommendations.append("图片质量整体良好，建议保持当前标准")

        recommendations.append("定期检查图片质量以确保演示效果")
        recommendations.append("考虑为不同使用场景准备不同分辨率的图片")

        return recommendations

    def _save_analysis_report(self, analysis_results: Dict):
        """保存分析报告"""
        # JSON格式报告
        json_path = os.path.join(self.output_dir, "image_quality_analysis.json")
        with open(json_path, 'w', encoding='utf-8') as f:
            json.dump(analysis_results, f, indent=2, ensure_ascii=False)

        # Markdown格式报告
        md_path = os.path.join(self.output_dir, "image_quality_report.md")
        self._generate_markdown_report(analysis_results, md_path)

        logger.info(f"分析报告已保存: {json_path}, {md_path}")

    def _generate_markdown_report(self, analysis_results: Dict, output_path: str):
        """生成Markdown格式报告"""
        with open(output_path, 'w', encoding='utf-8') as f:
            f.write("# PPT图片质量分析报告\n\n")

            # 基本信息
            file_info = analysis_results.get('file_info', {})
            f.write("## 基本信息\n")
            f.write(f"- **文件路径**: {file_info.get('path', 'N/A')}\n")
            f.write(f"- **总幻灯片数**: {file_info.get('total_slides', 0)}\n")
            f.write(f"- **分析时间**: {file_info.get('analysis_time', 'N/A')}\n\n")

            # 质量摘要
            quality_summary = analysis_results.get('quality_summary', {})
            if quality_summary:
                f.write("## 质量摘要\n")

                # 整体统计
                f.write(f"- **总图片数**: {quality_summary.get('total_images', 0)}\n")

                quality_stats = quality_summary.get('quality_statistics', {})
                if quality_stats:
                    f.write(f"- **平均质量分数**: {quality_stats.get('average_quality_score', 'N/A')}\n")
                    f.write(f"- **质量范围**: {quality_stats.get('min_quality_score', 'N/A')} - {quality_stats.get('max_quality_score', 'N/A')}\n")

                    quality_dist = quality_stats.get('quality_distribution', {})
                    f.write("- **质量分布**:\n")
                    f.write(f"  - 优秀 (≥0.8): {quality_dist.get('excellent', 0)}\n")
                    f.write(f"  - 良好 (0.6-0.8): {quality_dist.get('good', 0)}\n")
                    f.write(f"  - 一般 (0.4-0.6): {quality_dist.get('fair', 0)}\n")
                    f.write(f"  - 较差 (<0.4): {quality_dist.get('poor', 0)}\n")

                # 技术摘要
                tech_summary = quality_summary.get('technical_summary', {})
                if tech_summary:
                    f.write("\n### 技术指标\n")
                    f.write(f"- **平均亮度**: {tech_summary.get('average_brightness', 'N/A')}\n")
                    f.write(f"- **平均对比度**: {tech_summary.get('average_contrast', 'N/A')}\n")
                    f.write(f"- **平均锐度**: {tech_summary.get('average_sharpness', 'N/A')}\n")

                    if 'total_file_size' in tech_summary:
                        total_size_mb = tech_summary['total_file_size'] / (1024 * 1024)
                        avg_size_mb = tech_summary['average_file_size'] / (1024 * 1024)
                        f.write(f"- **总文件大小**: {total_size_mb:.1f} MB\n")
                        f.write(f"- **平均文件大小**: {avg_size_mb:.1f} MB\n")

                # 合规性摘要
                compliance_summary = quality_summary.get('compliance_summary', {})
                if compliance_summary:
                    f.write("\n### 合规性统计\n")
                    f.write(f"- **合规图片数**: {compliance_summary.get('compliant_images', 0)}\n")
                    f.write(f"- **合规率**: {compliance_summary.get('compliance_rate', 0):.1%}\n")
                    f.write(f"- **发现问题数**: {compliance_summary.get('total_issues', 0)}\n")
                    f.write(f"- **警告数**: {compliance_summary.get('total_warnings', 0)}\n")

            f.write("\n")

            # 详细分析
            image_analyses = analysis_results.get('image_analysis', [])
            if image_analyses:
                f.write("## 详细分析\n\n")

                for i, analysis in enumerate(image_analyses, 1):
                    location = analysis.get('location', {})
                    basic_info = analysis.get('basic_info', {})
                    quality_metrics = analysis.get('quality_metrics', {})

                    f.write(f"### 图片 {i}\n")
                    f.write(f"- **位置**: 幻灯片 {location.get('slide', 'N/A')}\n")
                    f.write(f"- **格式**: {basic_info.get('format', 'N/A')}\n")
                    f.write(f"- **尺寸**: {basic_info.get('width', 'N/A')}×{basic_info.get('height', 'N/A')}\n")
                    f.write(f"- **文件大小**: {basic_info.get('file_size_mb', 'N/A')} MB\n")

                    if quality_metrics:
                        f.write(f"- **质量分数**: {quality_metrics.get('overall_quality_score', 'N/A')}\n")
                        f.write(f"- **亮度**: {quality_metrics.get('brightness', 'N/A')}\n")
                        f.write(f"- **对比度**: {quality_metrics.get('contrast', 'N/A')}\n")
                        f.write(f"- **锐度**: {quality_metrics.get('sharpness', 'N/A')}\n")

                    compliance = analysis.get('compliance_check', {})
                    if compliance:
                        issues = compliance.get('issues', [])
                        warnings = compliance.get('warnings', [])

                        if issues:
                            f.write("- **问题**:\n")
                            for issue in issues:
                                f.write(f"  - {issue}\n")

                        if warnings:
                            f.write("- **警告**:\n")
                            for warning in warnings:
                                f.write(f"  - {warning}\n")

                    f.write("\n")

            # 优化建议
            recommendations = analysis_results.get('recommendations', [])
            if recommendations:
                f.write("## 优化建议\n\n")
                for i, rec in enumerate(recommendations, 1):
                    f.write(f"{i}. {rec}\n")

            # 技术问题
            technical_issues = analysis_results.get('technical_issues', [])
            if technical_issues:
                f.write("\n## 技术问题\n\n")
                for issue in technical_issues:
                    f.write(f"- 幻灯片 {issue.get('slide', 'N/A')}: {issue.get('issue', 'N/A')}\n")

            f.write("\n---\n")
            f.write(f"*报告生成时间: {time.strftime('%Y-%m-%d %H:%M:%S')}*\n")


def main():
    """主函数"""
    parser = argparse.ArgumentParser(description="AI-PPT-Assistant 图片质量分析器")
    parser.add_argument(
        'ppt_file',
        type=str,
        help='要分析的PPT文件路径'
    )
    parser.add_argument(
        '--output-dir',
        type=str,
        help='输出目录'
    )
    parser.add_argument(
        '--detailed',
        action='store_true',
        help='显示详细分析结果'
    )

    args = parser.parse_args()

    if not os.path.exists(args.ppt_file):
        print(f"错误: 文件不存在 {args.ppt_file}")
        sys.exit(1)

    try:
        analyzer = ImageQualityAnalyzer(output_dir=args.output_dir)
        results = analyzer.analyze_ppt_images(args.ppt_file)

        if 'error' in results:
            print(f"分析失败: {results['error']}")
            sys.exit(1)

        # 显示摘要
        quality_summary = results.get('quality_summary', {})
        print(f"\n=== PPT图片质量分析结果 ===")
        print(f"总图片数: {quality_summary.get('total_images', 0)}")

        quality_stats = quality_summary.get('quality_statistics', {})
        if quality_stats:
            print(f"平均质量分数: {quality_stats.get('average_quality_score', 'N/A')}")

            quality_dist = quality_stats.get('quality_distribution', {})
            print(f"质量分布: 优秀({quality_dist.get('excellent', 0)}) "
                  f"良好({quality_dist.get('good', 0)}) "
                  f"一般({quality_dist.get('fair', 0)}) "
                  f"较差({quality_dist.get('poor', 0)})")

        compliance_summary = quality_summary.get('compliance_summary', {})
        if compliance_summary:
            print(f"合规率: {compliance_summary.get('compliance_rate', 0):.1%}")
            print(f"发现问题: {compliance_summary.get('total_issues', 0)} 个")

        # 显示建议
        recommendations = results.get('recommendations', [])
        if recommendations:
            print(f"\n=== 优化建议 ===")
            for i, rec in enumerate(recommendations[:5], 1):  # 显示前5条
                print(f"{i}. {rec}")

        print(f"\n详细报告已保存到: {analyzer.output_dir}")

        # 详细模式
        if args.detailed:
            print(f"\n=== 详细分析结果 ===")
            print(json.dumps(results, indent=2, ensure_ascii=False))

    except Exception as e:
        logger.error(f"执行失败: {e}")
        sys.exit(1)


if __name__ == "__main__":
    main()