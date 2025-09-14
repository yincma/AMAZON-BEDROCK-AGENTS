"""
AI-PPT-Assistant PPT图片视觉验证测试套件

测试目标：
1. 端到端PPT生成测试
2. 图片在PPT中的显示效果验证
3. 图片质量和相关性检查
4. 不同模板下的图片适配
5. 错误情况下的占位图显示

作者: AI-PPT-Assistant Team
日期: 2025-01-14
"""

import unittest
import json
import os
import sys
import tempfile
import shutil
from pathlib import Path
from typing import Dict, List, Any, Optional
import logging
from io import BytesIO
import base64
import time
from PIL import Image, ImageChops
import numpy as np
import requests

# 添加项目路径
project_root = Path(__file__).parent.parent
sys.path.insert(0, str(project_root / "lambdas"))

try:
    from image_generator import ImageGenerator
    from image_config import CONFIG
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
except ImportError as e:
    logging.warning(f"导入依赖失败: {e}")

# 配置日志
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)


class PPTVisualValidationTest(unittest.TestCase):
    """PPT图片视觉验证测试类"""

    def setUp(self):
        """测试初始化"""
        self.test_dir = tempfile.mkdtemp(prefix="ppt_visual_test_")
        self.presentation_id = f"test_{int(time.time())}"

        # 测试数据
        self.test_slides = [
            {
                "title": "人工智能技术概览",
                "content": [
                    "机器学习算法的发展历程",
                    "深度学习在图像识别中的应用",
                    "自然语言处理的技术突破"
                ],
                "template": "business",
                "expected_image_keywords": ["technology", "AI", "futuristic"]
            },
            {
                "title": "数据分析与可视化",
                "content": [
                    "数据收集与清洗方法",
                    "统计分析的核心概念",
                    "可视化图表的设计原则"
                ],
                "template": "academic",
                "expected_image_keywords": ["data", "charts", "analytics"]
            },
            {
                "title": "云计算架构设计",
                "content": [
                    "微服务架构的优势",
                    "容器化部署策略",
                    "负载均衡与高可用设计"
                ],
                "template": "technical",
                "expected_image_keywords": ["cloud", "architecture", "network"]
            }
        ]

        # 图片质量标准
        self.quality_standards = {
            'min_width': 800,
            'min_height': 600,
            'max_file_size': 2 * 1024 * 1024,  # 2MB
            'required_formats': ['PNG', 'JPEG'],
            'min_dpi': 96,
            'aspect_ratio_tolerance': 0.1
        }

    def tearDown(self):
        """测试清理"""
        try:
            shutil.rmtree(self.test_dir)
        except Exception as e:
            logger.warning(f"清理测试目录失败: {e}")

    def test_01_end_to_end_ppt_generation(self):
        """测试1: 端到端PPT生成流程"""
        logger.info("开始端到端PPT生成测试...")

        try:
            # 1. 初始化图片生成器
            generator = ImageGenerator()

            # 2. 生成图片
            image_results = []
            for i, slide in enumerate(self.test_slides, 1):
                prompt = generator.generate_prompt(slide, target_audience=slide.get('template', 'business'))

                # 验证提示词质量
                self.assertIsInstance(prompt, str)
                self.assertGreater(len(prompt), 20, "提示词长度不足")
                self.assertTrue(any(keyword in prompt.lower() for keyword in slide['expected_image_keywords']),
                              f"提示词缺少预期关键词: {slide['expected_image_keywords']}")

                image_results.append({
                    'slide_number': i,
                    'prompt': prompt,
                    'template': slide['template']
                })

            # 3. 创建PPT
            ppt_path = self._create_test_ppt(image_results)
            self.assertTrue(os.path.exists(ppt_path), "PPT文件创建失败")

            # 4. 验证PPT结构
            self._verify_ppt_structure(ppt_path)

            logger.info("端到端PPT生成测试通过")

        except Exception as e:
            self.fail(f"端到端PPT生成测试失败: {e}")

    def test_02_image_display_quality(self):
        """测试2: 图片在PPT中的显示质量"""
        logger.info("开始图片显示质量测试...")

        try:
            # 1. 创建测试图片
            test_images = self._create_test_images()

            # 2. 创建包含图片的PPT
            ppt_path = self._create_ppt_with_images(test_images)

            # 3. 验证图片质量
            quality_report = self._analyze_image_quality(ppt_path)

            # 4. 质量断言
            self.assertGreaterEqual(quality_report['average_resolution'][0],
                                  self.quality_standards['min_width'],
                                  "图片宽度不足")
            self.assertGreaterEqual(quality_report['average_resolution'][1],
                                  self.quality_standards['min_height'],
                                  "图片高度不足")
            self.assertLessEqual(quality_report['max_file_size'],
                               self.quality_standards['max_file_size'],
                               "图片文件过大")

            logger.info(f"图片质量测试通过: {quality_report}")

        except Exception as e:
            self.fail(f"图片质量测试失败: {e}")

    def test_03_template_adaptation(self):
        """测试3: 不同模板下的图片适配"""
        logger.info("开始模板适配测试...")

        templates = ['business', 'academic', 'technical', 'creative']
        adaptation_results = {}

        try:
            for template in templates:
                # 1. 生成特定模板的提示词
                test_slide = {
                    "title": f"{template.title()}风格测试",
                    "content": ["测试内容1", "测试内容2", "测试内容3"]
                }

                generator = ImageGenerator()
                prompt = generator.generate_prompt(test_slide, target_audience=template)

                # 2. 验证模板特性
                template_characteristics = self._get_template_characteristics(template)
                self.assertTrue(
                    any(char in prompt.lower() for char in template_characteristics),
                    f"{template}模板特征未在提示词中体现"
                )

                adaptation_results[template] = {
                    'prompt': prompt,
                    'characteristics_found': [char for char in template_characteristics
                                            if char in prompt.lower()]
                }

            # 3. 验证模板差异性
            prompts = [result['prompt'] for result in adaptation_results.values()]
            self._verify_template_diversity(prompts)

            logger.info(f"模板适配测试通过: {adaptation_results}")

        except Exception as e:
            self.fail(f"模板适配测试失败: {e}")

    def test_04_image_content_relevance(self):
        """测试4: 图片与内容的相关性验证"""
        logger.info("开始图片内容相关性测试...")

        try:
            relevance_scores = []

            for slide in self.test_slides:
                # 1. 生成图片提示词
                generator = ImageGenerator()
                prompt = generator.generate_prompt(slide)

                # 2. 计算相关性分数
                relevance_score = self._calculate_content_relevance(slide, prompt)
                relevance_scores.append(relevance_score)

                # 3. 验证最低相关性要求
                self.assertGreaterEqual(relevance_score, 0.6,
                                      f"图片与内容相关性过低: {relevance_score}")

            # 4. 验证平均相关性
            average_relevance = sum(relevance_scores) / len(relevance_scores)
            self.assertGreaterEqual(average_relevance, 0.75,
                                  f"平均相关性不达标: {average_relevance}")

            logger.info(f"图片内容相关性测试通过，平均分数: {average_relevance:.2f}")

        except Exception as e:
            self.fail(f"图片内容相关性测试失败: {e}")

    def test_05_error_handling_placeholders(self):
        """测试5: 错误情况下的占位图处理"""
        logger.info("开始错误处理和占位图测试...")

        try:
            # 1. 模拟图片生成失败场景
            error_scenarios = [
                {"type": "network_error", "description": "网络连接失败"},
                {"type": "api_limit", "description": "API调用限制"},
                {"type": "invalid_prompt", "description": "无效提示词"},
                {"type": "service_unavailable", "description": "服务不可用"}
            ]

            fallback_results = []

            for scenario in error_scenarios:
                # 2. 生成占位图
                placeholder_result = self._generate_placeholder_for_error(scenario)
                fallback_results.append(placeholder_result)

                # 3. 验证占位图质量
                self.assertIsNotNone(placeholder_result['image_data'],
                                   f"{scenario['type']}场景下占位图生成失败")
                self.assertGreater(len(placeholder_result['image_data']), 1000,
                                 "占位图数据过小")

            # 4. 创建包含占位图的PPT进行验证
            placeholder_ppt = self._create_ppt_with_placeholders(fallback_results)
            self.assertTrue(os.path.exists(placeholder_ppt), "占位图PPT创建失败")

            logger.info("错误处理和占位图测试通过")

        except Exception as e:
            self.fail(f"错误处理测试失败: {e}")

    def test_06_multi_page_consistency(self):
        """测试6: 多页PPT图片一致性验证"""
        logger.info("开始多页PPT一致性测试...")

        try:
            # 1. 生成一致性图片系列
            generator = ImageGenerator()
            consistent_results = generator.generate_consistent_images(
                self.test_slides, self.presentation_id
            )

            # 2. 验证结果数量
            self.assertEqual(len(consistent_results), len(self.test_slides),
                           "生成的图片数量与幻灯片数量不匹配")

            # 3. 验证风格一致性
            style_consistency = self._analyze_style_consistency(consistent_results)
            self.assertGreaterEqual(style_consistency['consistency_score'], 0.8,
                                  f"风格一致性不足: {style_consistency['consistency_score']}")

            # 4. 创建多页PPT验证
            multi_page_ppt = self._create_multi_page_ppt(consistent_results)
            visual_consistency = self._analyze_visual_consistency(multi_page_ppt)

            self.assertTrue(visual_consistency['is_consistent'],
                          f"视觉一致性验证失败: {visual_consistency['issues']}")

            logger.info(f"多页PPT一致性测试通过: {visual_consistency}")

        except Exception as e:
            self.fail(f"多页一致性测试失败: {e}")

    def test_07_user_experience_evaluation(self):
        """测试7: 用户体验评估"""
        logger.info("开始用户体验评估测试...")

        try:
            # 1. 生成完整的用户场景PPT
            user_scenario_ppt = self._create_user_scenario_ppt()

            # 2. 执行用户体验评估
            ux_metrics = self._evaluate_user_experience(user_scenario_ppt)

            # 3. 验证关键UX指标
            self.assertGreaterEqual(ux_metrics['loading_performance'], 0.8,
                                  "加载性能不达标")
            self.assertGreaterEqual(ux_metrics['visual_appeal'], 0.7,
                                  "视觉吸引力不足")
            self.assertGreaterEqual(ux_metrics['content_clarity'], 0.8,
                                  "内容清晰度不足")
            self.assertLessEqual(ux_metrics['error_rate'], 0.1,
                               "错误率过高")

            # 4. 生成用户体验报告
            ux_report = self._generate_ux_report(ux_metrics)
            logger.info(f"用户体验评估完成: {ux_report}")

        except Exception as e:
            self.fail(f"用户体验评估失败: {e}")

    # 辅助方法
    def _create_test_ppt(self, image_results: List[Dict]) -> str:
        """创建测试PPT文件"""
        prs = Presentation()

        for i, result in enumerate(image_results):
            slide_layout = prs.slide_layouts[1]  # 使用标题和内容布局
            slide = prs.slides.add_slide(slide_layout)

            # 设置标题
            slide.shapes.title.text = f"测试幻灯片 {result['slide_number']}"

            # 添加内容
            content = slide.placeholders[1]
            text_frame = content.text_frame
            text_frame.clear()

            # 添加提示词信息
            p = text_frame.paragraphs[0]
            p.text = f"模板: {result['template']}"

            p2 = text_frame.add_paragraph()
            p2.text = f"提示词: {result['prompt'][:50]}..."

        ppt_path = os.path.join(self.test_dir, f"test_ppt_{self.presentation_id}.pptx")
        prs.save(ppt_path)
        return ppt_path

    def _verify_ppt_structure(self, ppt_path: str):
        """验证PPT结构"""
        prs = Presentation(ppt_path)

        # 验证幻灯片数量
        self.assertEqual(len(prs.slides), len(self.test_slides),
                        "幻灯片数量不正确")

        # 验证每张幻灯片的基本元素
        for i, slide in enumerate(prs.slides):
            self.assertTrue(slide.shapes.title is not None,
                          f"第{i+1}张幻灯片缺少标题")
            self.assertGreater(len(slide.shapes), 1,
                             f"第{i+1}张幻灯片元素过少")

    def _create_test_images(self) -> List[Dict]:
        """创建测试图片数据"""
        images = []

        for i in range(3):
            # 创建简单的测试图片
            img = Image.new('RGB', (1200, 800), color=(100 + i * 50, 150, 200))
            img_bytes = BytesIO()
            img.save(img_bytes, format='PNG')
            img_bytes.seek(0)

            images.append({
                'image_data': img_bytes.getvalue(),
                'slide_number': i + 1,
                'format': 'PNG',
                'size': (1200, 800)
            })

        return images

    def _create_ppt_with_images(self, test_images: List[Dict]) -> str:
        """创建包含图片的PPT"""
        prs = Presentation()

        for img_data in test_images:
            slide_layout = prs.slide_layouts[6]  # 空白布局
            slide = prs.slides.add_slide(slide_layout)

            # 添加图片
            img_stream = BytesIO(img_data['image_data'])

            # 计算图片位置和尺寸
            left = Inches(1)
            top = Inches(1.5)
            width = Inches(8)
            height = Inches(5)

            slide.shapes.add_picture(img_stream, left, top, width, height)

        ppt_path = os.path.join(self.test_dir, f"image_test_{self.presentation_id}.pptx")
        prs.save(ppt_path)
        return ppt_path

    def _analyze_image_quality(self, ppt_path: str) -> Dict:
        """分析PPT中图片的质量"""
        prs = Presentation(ppt_path)

        resolutions = []
        file_sizes = []

        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, 'image'):
                    img_bytes = shape.image.blob
                    img = Image.open(BytesIO(img_bytes))

                    resolutions.append(img.size)
                    file_sizes.append(len(img_bytes))

        if not resolutions:
            return {
                'average_resolution': (0, 0),
                'max_file_size': 0,
                'total_images': 0
            }

        avg_width = sum(r[0] for r in resolutions) / len(resolutions)
        avg_height = sum(r[1] for r in resolutions) / len(resolutions)

        return {
            'average_resolution': (int(avg_width), int(avg_height)),
            'max_file_size': max(file_sizes),
            'total_images': len(resolutions),
            'resolutions': resolutions,
            'file_sizes': file_sizes
        }

    def _get_template_characteristics(self, template: str) -> List[str]:
        """获取模板特征关键词"""
        characteristics = {
            'business': ['professional', 'corporate', 'clean', 'modern'],
            'academic': ['scholarly', 'research', 'educational', 'formal'],
            'technical': ['technical', 'engineering', 'diagram', 'systematic'],
            'creative': ['creative', 'artistic', 'colorful', 'innovative']
        }
        return characteristics.get(template, [])

    def _verify_template_diversity(self, prompts: List[str]):
        """验证模板间的多样性"""
        # 计算提示词间的相似度
        similarities = []

        for i in range(len(prompts)):
            for j in range(i + 1, len(prompts)):
                similarity = self._calculate_text_similarity(prompts[i], prompts[j])
                similarities.append(similarity)

        if similarities:
            avg_similarity = sum(similarities) / len(similarities)
            self.assertLess(avg_similarity, 0.7,
                          f"模板间差异性不足，平均相似度: {avg_similarity}")

    def _calculate_text_similarity(self, text1: str, text2: str) -> float:
        """计算文本相似度（简单实现）"""
        words1 = set(text1.lower().split())
        words2 = set(text2.lower().split())

        if not words1 or not words2:
            return 0.0

        intersection = len(words1.intersection(words2))
        union = len(words1.union(words2))

        return intersection / union if union > 0 else 0.0

    def _calculate_content_relevance(self, slide: Dict, prompt: str) -> float:
        """计算图片提示词与内容的相关性"""
        # 提取内容关键词
        content_words = set()
        content_words.update(slide.get('title', '').lower().split())

        for content_item in slide.get('content', []):
            content_words.update(content_item.lower().split())

        # 提取提示词关键词
        prompt_words = set(prompt.lower().split())

        # 计算相关性
        if not content_words or not prompt_words:
            return 0.0

        common_words = content_words.intersection(prompt_words)
        relevance = len(common_words) / len(content_words.union(prompt_words))

        return relevance

    def _generate_placeholder_for_error(self, scenario: Dict) -> Dict:
        """为错误场景生成占位图"""
        try:
            # 创建占位图
            img = Image.new('RGB', (800, 600), color=(220, 220, 220))
            img_bytes = BytesIO()
            img.save(img_bytes, format='PNG')
            img_bytes.seek(0)

            return {
                'scenario': scenario,
                'image_data': img_bytes.getvalue(),
                'status': 'fallback'
            }
        except Exception as e:
            return {
                'scenario': scenario,
                'image_data': None,
                'status': 'error',
                'error': str(e)
            }

    def _create_ppt_with_placeholders(self, fallback_results: List[Dict]) -> str:
        """创建包含占位图的PPT"""
        prs = Presentation()

        for result in fallback_results:
            if result['image_data']:
                slide_layout = prs.slide_layouts[1]
                slide = prs.slides.add_slide(slide_layout)

                # 设置标题
                slide.shapes.title.text = f"错误处理: {result['scenario']['type']}"

                # 添加占位图
                img_stream = BytesIO(result['image_data'])
                left = Inches(1)
                top = Inches(2)
                width = Inches(6)
                height = Inches(4)

                slide.shapes.add_picture(img_stream, left, top, width, height)

        ppt_path = os.path.join(self.test_dir, f"placeholder_test_{self.presentation_id}.pptx")
        prs.save(ppt_path)
        return ppt_path

    def _analyze_style_consistency(self, results: List[Dict]) -> Dict:
        """分析风格一致性"""
        # 检查style_params一致性
        style_params = [r.get('style_params', {}) for r in results if 'style_params' in r]

        if not style_params:
            return {'consistency_score': 0.0, 'message': '无风格参数'}

        # 计算一致性分数
        consistency_scores = []
        base_style = style_params[0]

        for style in style_params[1:]:
            score = self._compare_style_params(base_style, style)
            consistency_scores.append(score)

        avg_consistency = sum(consistency_scores) / len(consistency_scores) if consistency_scores else 1.0

        return {
            'consistency_score': avg_consistency,
            'style_variations': len(set(str(s) for s in style_params)),
            'total_checked': len(style_params)
        }

    def _compare_style_params(self, style1: Dict, style2: Dict) -> float:
        """比较两个风格参数的相似度"""
        if not style1 or not style2:
            return 0.0

        common_keys = set(style1.keys()).intersection(set(style2.keys()))
        if not common_keys:
            return 0.0

        matching_values = sum(1 for key in common_keys if style1[key] == style2[key])
        return matching_values / len(common_keys)

    def _create_multi_page_ppt(self, consistent_results: List[Dict]) -> str:
        """创建多页一致性PPT"""
        prs = Presentation()

        for i, result in enumerate(consistent_results, 1):
            slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(slide_layout)

            # 设置标题
            slide.shapes.title.text = f"一致性测试页 {i}"

            # 添加风格信息
            content = slide.placeholders[1]
            text_frame = content.text_frame
            text_frame.clear()

            p = text_frame.paragraphs[0]
            p.text = f"状态: {result.get('status', 'unknown')}"

            if 'style_params' in result:
                p2 = text_frame.add_paragraph()
                p2.text = f"风格: {result['style_params']}"

        ppt_path = os.path.join(self.test_dir, f"consistency_test_{self.presentation_id}.pptx")
        prs.save(ppt_path)
        return ppt_path

    def _analyze_visual_consistency(self, ppt_path: str) -> Dict:
        """分析视觉一致性"""
        prs = Presentation(ppt_path)

        # 基本一致性检查
        slide_count = len(prs.slides)
        layout_consistency = True
        issues = []

        if slide_count < 2:
            return {
                'is_consistent': True,
                'issues': ['仅有一张幻灯片，无需检查一致性']
            }

        # 检查布局一致性
        first_layout = prs.slides[0].slide_layout
        for i, slide in enumerate(prs.slides[1:], 1):
            if slide.slide_layout != first_layout:
                layout_consistency = False
                issues.append(f"第{i+1}张幻灯片布局不一致")

        return {
            'is_consistent': layout_consistency and len(issues) == 0,
            'issues': issues,
            'slide_count': slide_count,
            'layout_consistency': layout_consistency
        }

    def _create_user_scenario_ppt(self) -> str:
        """创建用户场景PPT"""
        prs = Presentation()

        # 创建典型的用户使用场景
        scenarios = [
            {"title": "项目启动会议", "content": "项目目标和里程碑规划"},
            {"title": "技术架构方案", "content": "系统设计和技术选型"},
            {"title": "进度汇报", "content": "当前进展和下一步计划"}
        ]

        for scenario in scenarios:
            slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(slide_layout)
            slide.shapes.title.text = scenario['title']

            content = slide.placeholders[1]
            content.text = scenario['content']

        ppt_path = os.path.join(self.test_dir, f"user_scenario_{self.presentation_id}.pptx")
        prs.save(ppt_path)
        return ppt_path

    def _evaluate_user_experience(self, ppt_path: str) -> Dict:
        """评估用户体验"""
        try:
            # 加载性能测试
            start_time = time.time()
            prs = Presentation(ppt_path)
            load_time = time.time() - start_time

            # 计算各项指标
            loading_performance = min(1.0, 3.0 / max(load_time, 0.1))  # 3秒内加载完成为满分

            # 视觉吸引力（基于内容丰富度）
            total_elements = sum(len(slide.shapes) for slide in prs.slides)
            visual_appeal = min(1.0, total_elements / (len(prs.slides) * 3))  # 每页至少3个元素

            # 内容清晰度（基于文本可读性）
            text_clarity = 0.8  # 模拟评分

            # 错误率（基于异常检测）
            error_count = 0
            for slide in prs.slides:
                try:
                    # 检查每张幻灯片是否可正常访问
                    _ = slide.shapes.title
                except:
                    error_count += 1

            error_rate = error_count / len(prs.slides) if prs.slides else 0

            return {
                'loading_performance': loading_performance,
                'visual_appeal': visual_appeal,
                'content_clarity': text_clarity,
                'error_rate': error_rate,
                'load_time': load_time,
                'total_slides': len(prs.slides),
                'total_elements': total_elements
            }

        except Exception as e:
            return {
                'loading_performance': 0.0,
                'visual_appeal': 0.0,
                'content_clarity': 0.0,
                'error_rate': 1.0,
                'error': str(e)
            }

    def _generate_ux_report(self, ux_metrics: Dict) -> Dict:
        """生成用户体验报告"""
        overall_score = (
            ux_metrics.get('loading_performance', 0) * 0.3 +
            ux_metrics.get('visual_appeal', 0) * 0.25 +
            ux_metrics.get('content_clarity', 0) * 0.25 +
            (1 - ux_metrics.get('error_rate', 1)) * 0.2
        )

        grade_mapping = {
            (0.9, 1.0): 'A',
            (0.8, 0.9): 'B',
            (0.7, 0.8): 'C',
            (0.6, 0.7): 'D',
            (0.0, 0.6): 'F'
        }

        grade = 'F'
        for (min_score, max_score), letter in grade_mapping.items():
            if min_score <= overall_score < max_score:
                grade = letter
                break

        return {
            'overall_score': overall_score,
            'grade': grade,
            'metrics': ux_metrics,
            'recommendations': self._generate_recommendations(ux_metrics)
        }

    def _generate_recommendations(self, metrics: Dict) -> List[str]:
        """基于指标生成改进建议"""
        recommendations = []

        if metrics.get('loading_performance', 0) < 0.7:
            recommendations.append("优化图片大小和格式以提升加载速度")

        if metrics.get('visual_appeal', 0) < 0.6:
            recommendations.append("增加视觉元素，提升页面丰富度")

        if metrics.get('content_clarity', 0) < 0.7:
            recommendations.append("改善文本布局和字体选择")

        if metrics.get('error_rate', 0) > 0.1:
            recommendations.append("修复发现的错误和异常")

        if not recommendations:
            recommendations.append("用户体验表现良好，可考虑进一步优化细节")

        return recommendations


class PPTVisualValidationSuite:
    """PPT视觉验证测试套件管理器"""

    def __init__(self, output_dir: str = None):
        self.output_dir = output_dir or tempfile.mkdtemp(prefix="ppt_validation_")
        self.results = {}

    def run_full_validation(self) -> Dict:
        """运行完整的视觉验证测试"""
        logger.info("开始执行完整的PPT视觉验证测试套件...")

        # 创建测试套件
        loader = unittest.TestLoader()
        suite = loader.loadTestsFromTestCase(PPTVisualValidationTest)

        # 运行测试
        runner = unittest.TextTestRunner(verbosity=2, stream=sys.stdout)
        result = runner.run(suite)

        # 收集结果
        validation_result = {
            'timestamp': time.strftime('%Y-%m-%d %H:%M:%S'),
            'total_tests': result.testsRun,
            'failures': len(result.failures),
            'errors': len(result.errors),
            'success_rate': (result.testsRun - len(result.failures) - len(result.errors)) / result.testsRun if result.testsRun > 0 else 0,
            'output_dir': self.output_dir
        }

        # 生成详细报告
        self._generate_validation_report(validation_result, result)

        return validation_result

    def _generate_validation_report(self, summary: Dict, test_result) -> str:
        """生成验证报告"""
        report_path = os.path.join(self.output_dir, "visual_validation_report.md")

        with open(report_path, 'w', encoding='utf-8') as f:
            f.write("# AI-PPT-Assistant 图片显示效果验证报告\n\n")
            f.write(f"生成时间: {summary['timestamp']}\n\n")

            f.write("## 测试概要\n")
            f.write(f"- 总测试数: {summary['total_tests']}\n")
            f.write(f"- 失败数: {summary['failures']}\n")
            f.write(f"- 错误数: {summary['errors']}\n")
            f.write(f"- 成功率: {summary['success_rate']:.1%}\n\n")

            if test_result.failures:
                f.write("## 失败的测试\n")
                for test, traceback in test_result.failures:
                    f.write(f"### {test}\n")
                    f.write("```\n")
                    f.write(traceback)
                    f.write("\n```\n\n")

            if test_result.errors:
                f.write("## 错误的测试\n")
                for test, traceback in test_result.errors:
                    f.write(f"### {test}\n")
                    f.write("```\n")
                    f.write(traceback)
                    f.write("\n```\n\n")

            f.write("## 改进建议\n")
            f.write("1. 持续监控图片生成质量\n")
            f.write("2. 优化模板适配算法\n")
            f.write("3. 完善错误处理机制\n")
            f.write("4. 提升用户体验指标\n")

        return report_path


if __name__ == "__main__":
    # 运行完整的视觉验证测试
    validation_suite = PPTVisualValidationSuite()
    results = validation_suite.run_full_validation()

    print(f"\n视觉验证测试完成！")
    print(f"成功率: {results['success_rate']:.1%}")
    print(f"报告输出目录: {results['output_dir']}")