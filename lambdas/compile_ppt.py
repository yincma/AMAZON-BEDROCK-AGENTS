"""
PPT编译Lambda函数 - 处理PPT编译请求
"""

import json
import os
import logging
import sys

# 添加src路径以导入我们的模块
sys.path.insert(0, '/opt/python/lib/python3.13/site-packages')
sys.path.insert(0, os.path.join(os.path.dirname(__file__), '../src'))

try:
    from src.ppt_compiler import PPTCompiler, create_pptx_from_content
    from src.ppt_validator import PPTXValidator
    from src.file_utils import ensure_tmp_space, clean_tmp_files
except ImportError:
    # 如果src模块不可用，使用本地实现
    print("Warning: src modules not available, using fallback implementation")

# 配置日志
logger = logging.getLogger()
logger.setLevel(logging.INFO)


def lambda_handler(event, context):
    """
    PPT编译Lambda处理函数

    Args:
        event: Lambda事件，包含presentation_id
        context: Lambda上下文

    Returns:
        dict: 响应包含状态码、下载URL等
    """
    try:
        logger.info(f"PPT compilation request received: {json.dumps(event)}")

        # 1. 解析请求
        presentation_id = extract_presentation_id(event)
        if not presentation_id:
            return format_error_response(400, "presentation_id is required")

        # 2. 检查临时空间
        try:
            ensure_tmp_space(min_free_mb=50)  # 确保至少有50MB空间
        except Exception as e:
            logger.warning(f"Failed to ensure tmp space: {e}")
            # 尝试清理并继续
            clean_tmp_files()

        # 3. 初始化编译器
        compiler = PPTCompiler()

        # 4. 编译PPT
        download_url = compiler.compile_ppt(presentation_id)

        # 5. 返回成功响应
        response = {
            'statusCode': 200,
            'headers': {
                'Content-Type': 'application/json',
                'Access-Control-Allow-Origin': '*',
                'Access-Control-Allow-Headers': 'Content-Type,X-Api-Key,Accept',
                'Access-Control-Allow-Methods': 'GET,POST,OPTIONS'
            },
            'body': json.dumps({
                'presentation_id': presentation_id,
                'download_url': download_url,
                'status': 'completed',
                'message': 'PPT compiled successfully'
            })
        }

        logger.info(f"PPT compilation completed successfully for {presentation_id}")
        return response

    except Exception as e:
        logger.error(f"PPT compilation failed: {str(e)}", exc_info=True)
        return format_error_response(500, f"Compilation failed: {str(e)}")

    finally:
        # 清理临时文件
        try:
            clean_tmp_files(max_age_hours=0.5)  # 清理超过30分钟的文件
        except Exception as e:
            logger.warning(f"Failed to clean tmp files: {e}")


def extract_presentation_id(event):
    """
    从Lambda事件中提取presentation_id

    Args:
        event: Lambda事件

    Returns:
        str: presentation_id或None
    """
    # 1. 直接从event中获取
    if 'presentation_id' in event:
        return event['presentation_id']

    # 2. 从body中解析
    if 'body' in event:
        try:
            if isinstance(event['body'], str):
                body = json.loads(event['body'])
            else:
                body = event['body']

            if 'presentation_id' in body:
                return body['presentation_id']
        except (json.JSONDecodeError, TypeError):
            logger.warning("Failed to parse body as JSON")

    # 3. 从路径参数获取
    if 'pathParameters' in event and event['pathParameters']:
        if 'presentation_id' in event['pathParameters']:
            return event['pathParameters']['presentation_id']

    # 4. 从查询参数获取
    if 'queryStringParameters' in event and event['queryStringParameters']:
        if 'presentation_id' in event['queryStringParameters']:
            return event['queryStringParameters']['presentation_id']

    return None


def format_error_response(status_code, error_message):
    """
    创建错误响应

    Args:
        status_code: HTTP状态码
        error_message: 错误消息

    Returns:
        dict: 错误响应
    """
    return {
        'statusCode': status_code,
        'headers': {
            'Content-Type': 'application/json',
            'Access-Control-Allow-Origin': '*',
            'Access-Control-Allow-Headers': 'Content-Type,X-Amz-Date,Authorization,X-Api-Key,X-Amz-Security-Token',
            'Access-Control-Allow-Methods': 'POST,OPTIONS'
        },
        'body': json.dumps({
            'error': error_message,
            'status': 'failed'
        })
    }


# 简化的编译函数（如果主模块不可用）
def simple_compile_ppt(presentation_id):
    """
    简化的PPT编译实现（作为后备）

    Args:
        presentation_id: 演示文稿ID

    Returns:
        str: 下载URL
    """
    import boto3
    from pptx import Presentation
    from io import BytesIO
    import tempfile

    s3_client = boto3.client('s3')
    bucket_name = os.environ.get('S3_BUCKET', 'ai-ppt-presentations-dev')

    try:
        # 读取内容
        content_key = f"presentations/{presentation_id}/content.json"
        response = s3_client.get_object(Bucket=bucket_name, Key=content_key)
        content = json.loads(response['Body'].read().decode('utf-8'))

        # 创建PPT
        prs = Presentation()

        if 'slides' in content:
            for slide_data in content['slides']:
                slide_layout = prs.slide_layouts[1]
                slide = prs.slides.add_slide(slide_layout)

                # 添加标题
                if 'title' in slide_data:
                    title = slide.shapes.title
                    title.text = slide_data['title']

                # 添加内容
                if 'bullet_points' in slide_data:
                    content_placeholder = slide.placeholders[1]
                    text_frame = content_placeholder.text_frame
                    text_frame.clear()

                    for i, point in enumerate(slide_data['bullet_points']):
                        if i == 0:
                            p = text_frame.paragraphs[0]
                        else:
                            p = text_frame.add_paragraph()
                        p.text = point

        # 保存到字节流
        pptx_bytes = BytesIO()
        prs.save(pptx_bytes)
        pptx_bytes.seek(0)

        # 上传到S3
        s3_key = f"presentations/{presentation_id}/output/presentation.pptx"
        s3_client.put_object(
            Bucket=bucket_name,
            Key=s3_key,
            Body=pptx_bytes.getvalue(),
            ContentType='application/vnd.openxmlformats-officedocument.presentationml.presentation'
        )

        # 生成下载链接
        download_url = s3_client.generate_presigned_url(
            'get_object',
            Params={'Bucket': bucket_name, 'Key': s3_key},
            ExpiresIn=3600
        )

        return download_url

    except Exception as e:
        logger.error(f"Simple compilation failed: {e}")
        raise


def lambda_handler(event, context):
    """
    Lambda入口点 - 兼容不同的调用方式
    """
    return handler(event, context)


# 测试函数
if __name__ == "__main__":
    # 本地测试
    test_event = {
        'body': json.dumps({
            'presentation_id': 'test-123'
        })
    }

    test_context = type('obj', (object,), {
        'function_name': 'test_compile_ppt',
        'function_version': '$LATEST',
        'memory_limit_in_mb': 1024
    })

    result = handler(test_event, test_context)
    print(json.dumps(result, indent=2))