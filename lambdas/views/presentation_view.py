"""
Presentation View - Rendering Layer for PPTX Generation
Handles PowerPoint file creation, layouts, and styling using python-pptx
"""

import io
import os
import sys
from enum import Enum
from typing import Any, Dict, List, Optional

from aws_lambda_powertools import Logger, Metrics, Tracer
from aws_lambda_powertools.metrics import MetricUnit
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# Add parent directory to path for imports
sys.path.append(os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

# Import styles
from views.presentation_styles import ColorScheme, SlideStyle

# Import interfaces
from interfaces.presentation_view_interface import (
    IPresentationView,
)

# Initialize AWS Lambda Powertools
logger = Logger()
tracer = Tracer()
metrics = Metrics()


# Layout types
class SlideLayout(Enum):
    TITLE = 0  # Title slide
    CONTENT = 1  # Title and content
    SECTION = 2  # Section header
    TWO_CONTENT = 3  # Two content columns
    COMPARISON = 4  # Side-by-side comparison
    TITLE_ONLY = 5  # Title only
    BLANK = 6  # Blank slide
    CONTENT_CAPTION = 7  # Content with caption
    PICTURE_CAPTION = 8  # Picture with caption


# Styles are now imported from views.presentation_styles


class PresentationView(IPresentationView):
    """View layer for PowerPoint presentation generation with interface compliance"""

    def __init__(self, template_path: Optional[str] = None):
        """
        Initialize presentation view

        Args:
            template_path: Optional path to PPTX template file
        """
        if template_path and os.path.exists(template_path):
            self.presentation = Presentation(template_path)
            logger.info(f"Loaded template: {template_path}")
        else:
            self.presentation = Presentation()
            logger.info("Created blank presentation")

        self.style = SlideStyle()
        self._setup_default_layouts()

    def _setup_default_layouts(self):
        """Setup default slide layouts if not using template"""
        # This would be called when creating from scratch
        # python-pptx uses the default template layouts

    @tracer.capture_method
    def set_style(self, style: SlideStyle):
        """Set the presentation style"""
        self.style = style
        logger.info(f"Set presentation style: {style.title_font}")

    @tracer.capture_method
    def add_title_slide(
        self,
        title: str,
        subtitle: Optional[str] = None,
        author: Optional[str] = None,
        date: Optional[str] = None,
    ) -> Any:
        """Add a title slide"""
        try:
            slide_layout = self.presentation.slide_layouts[SlideLayout.TITLE.value]
            slide = self.presentation.slides.add_slide(slide_layout)

            # Set title
            title_shape = slide.shapes.title
            if title_shape:
                title_shape.text = title
                self._apply_text_style(
                    title_shape.text_frame,
                    self.style.title_font,
                    self.style.title_size,
                    self.style.color_scheme["primary"],
                    bold=True,
                )

            # Set subtitle
            if subtitle and len(slide.placeholders) > 1:
                subtitle_shape = slide.placeholders[1]
                subtitle_text = subtitle
                if author:
                    subtitle_text += f"\n{author}"
                if date:
                    subtitle_text += f"\n{date}"
                subtitle_shape.text = subtitle_text
                self._apply_text_style(
                    subtitle_shape.text_frame,
                    self.style.subtitle_font,
                    self.style.subtitle_size,
                    self.style.color_scheme["secondary"],
                )

            logger.info(f"Added title slide: {title}")
            return slide

        except Exception as e:
            logger.error(f"Error adding title slide: {str(e)}")
            raise

    @tracer.capture_method
    def add_content_slide(
        self,
        title: str,
        content: List[str],
        level: int = 0,
        slide_number: Optional[int] = None,
    ) -> Any:
        """Add a content slide with bullet points"""
        try:
            slide_layout = self.presentation.slide_layouts[SlideLayout.CONTENT.value]
            slide = self.presentation.slides.add_slide(slide_layout)

            # Set title
            title_shape = slide.shapes.title
            if title_shape:
                display_title = title
                if slide_number:
                    display_title = f"{slide_number}. {title}"
                title_shape.text = display_title
                self._apply_text_style(
                    title_shape.text_frame,
                    self.style.title_font,
                    self.style.title_size - 8,  # Slightly smaller for content slides
                    self.style.color_scheme["primary"],
                    bold=True,
                )

            # Add content
            content_shape = None
            for shape in slide.shapes:
                if shape.has_text_frame and shape != title_shape:
                    content_shape = shape
                    break

            if content_shape:
                text_frame = content_shape.text_frame
                text_frame.clear()  # Clear existing text

                for i, item in enumerate(content):
                    if i == 0:
                        paragraph = text_frame.paragraphs[0]
                    else:
                        paragraph = text_frame.add_paragraph()

                    paragraph.text = item
                    paragraph.level = level
                    self._apply_bullet_style(
                        paragraph,
                        self.style.bullet_font,
                        self.style.bullet_size,
                        self.style.color_scheme["text"],
                    )

            logger.info(f"Added content slide: {title} with {len(content)} items")
            return slide

        except Exception as e:
            logger.error(f"Error adding content slide: {str(e)}")
            raise

    @tracer.capture_method
    def add_two_content_slide(
        self,
        title: str,
        left_content: List[str],
        right_content: List[str],
        left_title: Optional[str] = None,
        right_title: Optional[str] = None,
    ) -> Any:
        """Add a slide with two content columns"""
        try:
            slide_layout = self.presentation.slide_layouts[
                SlideLayout.TWO_CONTENT.value
            ]
            slide = self.presentation.slides.add_slide(slide_layout)

            # Set title
            title_shape = slide.shapes.title
            if title_shape:
                title_shape.text = title
                self._apply_text_style(
                    title_shape.text_frame,
                    self.style.title_font,
                    self.style.title_size - 8,
                    self.style.color_scheme["primary"],
                    bold=True,
                )

            # Find content placeholders
            placeholders = [
                shape
                for shape in slide.shapes
                if shape.has_text_frame and shape != title_shape
            ]

            # Add left content
            if len(placeholders) > 0:
                self._add_content_to_shape(placeholders[0], left_content, left_title)

            # Add right content
            if len(placeholders) > 1:
                self._add_content_to_shape(placeholders[1], right_content, right_title)

            logger.info(f"Added two-content slide: {title}")
            return slide

        except Exception as e:
            logger.error(f"Error adding two-content slide: {str(e)}")
            raise

    @tracer.capture_method
    def add_image_slide(
        self,
        title: str,
        image_data: bytes,
        caption: Optional[str] = None,
        layout: str = "center",
    ) -> Any:
        """Add a slide with an image"""
        try:
            slide_layout = self.presentation.slide_layouts[
                SlideLayout.PICTURE_CAPTION.value
            ]
            slide = self.presentation.slides.add_slide(slide_layout)

            # Set title
            title_shape = slide.shapes.title
            if title_shape:
                title_shape.text = title
                self._apply_text_style(
                    title_shape.text_frame,
                    self.style.title_font,
                    self.style.title_size - 8,
                    self.style.color_scheme["primary"],
                    bold=True,
                )

            # Add image
            image_stream = io.BytesIO(image_data)

            if layout == "center":
                left = Inches(1.5)
                top = Inches(2)
                width = Inches(7)
                height = Inches(4)
            elif layout == "left":
                left = Inches(0.5)
                top = Inches(2)
                width = Inches(4.5)
                height = Inches(4)
            else:  # right
                left = Inches(5)
                top = Inches(2)
                width = Inches(4.5)
                height = Inches(4)

            slide.shapes.add_picture(image_stream, left, top, width, height)

            # Add caption if provided
            if caption:
                caption_left = left
                caption_top = top + height + Inches(0.2)
                caption_width = width
                caption_height = Inches(0.5)

                textbox = slide.shapes.add_textbox(
                    caption_left, caption_top, caption_width, caption_height
                )
                text_frame = textbox.text_frame
                text_frame.text = caption
                text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                self._apply_text_style(
                    text_frame,
                    self.style.content_font,
                    self.style.content_size - 2,
                    self.style.color_scheme["secondary"],
                    italic=True,
                )

            logger.info(f"Added image slide: {title}")
            return slide

        except Exception as e:
            logger.error(f"Error adding image slide: {str(e)}")
            raise

    @tracer.capture_method
    def add_chart_slide(
        self,
        title: str,
        chart_type: str,
        data: Dict[str, Any],
        chart_title: Optional[str] = None,
    ) -> Any:
        """Add a slide with a chart"""
        try:
            slide_layout = self.presentation.slide_layouts[SlideLayout.CONTENT.value]
            slide = self.presentation.slides.add_slide(slide_layout)

            # Set title
            title_shape = slide.shapes.title
            if title_shape:
                title_shape.text = title
                self._apply_text_style(
                    title_shape.text_frame,
                    self.style.title_font,
                    self.style.title_size - 8,
                    self.style.color_scheme["primary"],
                    bold=True,
                )

            # Define chart position and size
            x, y, cx, cy = Inches(1), Inches(2), Inches(8), Inches(4.5)

            # Create chart data
            chart_data = CategoryChartData()

            # Process data based on structure
            if "categories" in data and "series" in data:
                chart_data.categories = data["categories"]
                for series_name, series_values in data["series"].items():
                    chart_data.add_series(series_name, series_values)
            else:
                # Simple data structure
                chart_data.categories = list(data.keys())
                chart_data.add_series("Values", list(data.values()))

            # Determine chart type
            chart_type_mapping = {
                "column": XL_CHART_TYPE.COLUMN_CLUSTERED,
                "bar": XL_CHART_TYPE.BAR_CLUSTERED,
                "line": XL_CHART_TYPE.LINE,
                "pie": XL_CHART_TYPE.PIE,
                "area": XL_CHART_TYPE.AREA,
            }

            xl_chart_type = chart_type_mapping.get(
                chart_type.lower(), XL_CHART_TYPE.COLUMN_CLUSTERED
            )

            # Add chart
            chart = slide.shapes.add_chart(
                xl_chart_type, x, y, cx, cy, chart_data
            ).chart

            # Set chart title
            if chart_title:
                chart.has_title = True
                chart.chart_title.text_frame.text = chart_title

            logger.info(f"Added chart slide: {title} with {chart_type} chart")
            return slide

        except Exception as e:
            logger.error(f"Error adding chart slide: {str(e)}")
            raise

    @tracer.capture_method
    def add_section_slide(self, title: str, subtitle: Optional[str] = None) -> Any:
        """Add a section divider slide"""
        try:
            slide_layout = self.presentation.slide_layouts[SlideLayout.SECTION.value]
            slide = self.presentation.slides.add_slide(slide_layout)

            # Set title
            title_shape = slide.shapes.title
            if title_shape:
                title_shape.text = title
                self._apply_text_style(
                    title_shape.text_frame,
                    self.style.title_font,
                    self.style.title_size,
                    self.style.color_scheme["primary"],
                    bold=True,
                )

            # Set subtitle if provided
            if subtitle:
                for shape in slide.shapes:
                    if shape.has_text_frame and shape != title_shape:
                        shape.text = subtitle
                        self._apply_text_style(
                            shape.text_frame,
                            self.style.subtitle_font,
                            self.style.subtitle_size - 4,
                            self.style.color_scheme["secondary"],
                        )
                        break

            logger.info(f"Added section slide: {title}")
            return slide

        except Exception as e:
            logger.error(f"Error adding section slide: {str(e)}")
            raise

    @tracer.capture_method
    def add_closing_slide(
        self,
        title: str = "Thank You",
        contact_info: Optional[Dict[str, str]] = None,
        additional_text: Optional[str] = None,
    ) -> Any:
        """Add a closing/thank you slide"""
        try:
            slide_layout = self.presentation.slide_layouts[SlideLayout.TITLE_ONLY.value]
            slide = self.presentation.slides.add_slide(slide_layout)

            # Set title
            title_shape = slide.shapes.title
            if title_shape:
                title_shape.text = title
                self._apply_text_style(
                    title_shape.text_frame,
                    self.style.title_font,
                    self.style.title_size + 8,  # Larger for impact
                    self.style.color_scheme["primary"],
                    bold=True,
                )

            # Add contact info or additional text
            if contact_info or additional_text:
                left = Inches(2)
                top = Inches(4)
                width = Inches(6)
                height = Inches(2)

                textbox = slide.shapes.add_textbox(left, top, width, height)
                text_frame = textbox.text_frame
                text_frame.clear()

                if additional_text:
                    paragraph = text_frame.paragraphs[0]
                    paragraph.text = additional_text
                    paragraph.alignment = PP_ALIGN.CENTER
                    self._apply_text_style(
                        text_frame,
                        self.style.content_font,
                        self.style.content_size,
                        self.style.color_scheme["text"],
                    )

                if contact_info:
                    for key, value in contact_info.items():
                        paragraph = text_frame.add_paragraph()
                        paragraph.text = f"{key}: {value}"
                        paragraph.alignment = PP_ALIGN.CENTER
                        self._apply_text_style(
                            text_frame,
                            self.style.content_font,
                            self.style.content_size - 2,
                            self.style.color_scheme["secondary"],
                        )

            logger.info(f"Added closing slide: {title}")
            return slide

        except Exception as e:
            logger.error(f"Error adding closing slide: {str(e)}")
            raise

    @tracer.capture_method
    def add_speaker_notes(self, slide: Any, notes: str):
        """Add speaker notes to a slide"""
        try:
            notes_slide = slide.notes_slide
            text_frame = notes_slide.notes_text_frame
            text_frame.text = notes
            logger.info("Added speaker notes to slide")
        except Exception as e:
            logger.error(f"Error adding speaker notes: {str(e)}")

    @tracer.capture_method
    def save_presentation(self, output_path: Optional[str] = None) -> bytes:
        """
        Save the presentation to file or return as bytes

        Args:
            output_path: Optional file path to save presentation

        Returns:
            Presentation file as bytes
        """
        try:
            # Save to BytesIO
            output = io.BytesIO()
            self.presentation.save(output)
            output.seek(0)
            presentation_bytes = output.read()

            # Optionally save to file
            if output_path:
                with open(output_path, "wb") as f:
                    f.write(presentation_bytes)
                logger.info(f"Saved presentation to: {output_path}")

            metrics.add_metric(
                name="PresentationGenerated", unit=MetricUnit.Count, value=1
            )
            metrics.add_metric(
                name="SlideCount",
                unit=MetricUnit.Count,
                value=len(self.presentation.slides),
            )

            return presentation_bytes

        except Exception as e:
            logger.error(f"Error saving presentation: {str(e)}")
            raise

    # Helper methods
    def _apply_text_style(
        self,
        text_frame,
        font_name: str,
        font_size: int,
        color: RGBColor,
        bold: bool = False,
        italic: bool = False,
    ):
        """Apply text styling to a text frame"""
        for paragraph in text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.name = font_name
                run.font.size = Pt(font_size)
                run.font.color.rgb = color
                run.font.bold = bold
                run.font.italic = italic

    def _apply_bullet_style(
        self, paragraph, font_name: str, font_size: int, color: RGBColor
    ):
        """Apply bullet point styling"""
        paragraph.font.name = font_name
        paragraph.font.size = Pt(font_size)
        paragraph.font.color.rgb = color

    def _add_content_to_shape(
        self, shape, content: List[str], title: Optional[str] = None
    ):
        """Add content to a text shape"""
        text_frame = shape.text_frame
        text_frame.clear()

        # Add title if provided
        if title:
            paragraph = text_frame.paragraphs[0]
            paragraph.text = title
            paragraph.font.bold = True
            paragraph.font.size = Pt(self.style.content_size + 2)
            paragraph.font.color.rgb = self.style.color_scheme["primary"]
            text_frame.add_paragraph()  # Add spacing

        # Add content items
        for i, item in enumerate(content):
            if i == 0 and not title:
                paragraph = text_frame.paragraphs[0]
            else:
                paragraph = text_frame.add_paragraph()
            paragraph.text = item
            self._apply_bullet_style(
                paragraph,
                self.style.bullet_font,
                self.style.bullet_size,
                self.style.color_scheme["text"],
            )

    def get_slide_count(self) -> int:
        """Get the current number of slides"""
        return len(self.presentation.slides)

    def apply_template(self, template_data: bytes):
        """Apply a template to the presentation"""
        try:
            # Load template from bytes
            template_stream = io.BytesIO(template_data)
            template_presentation = Presentation(template_stream)

            # Copy slide layouts from template
            self.presentation = template_presentation
            logger.info("Applied template to presentation")

        except Exception as e:
            logger.error(f"Error applying template: {str(e)}")
            raise


# Helper function to create styled presentation
def create_styled_presentation(
    style_name: str = "professional", template_data: Optional[bytes] = None
) -> PresentationView:
    """
    Create a presentation with predefined style

    Args:
        style_name: Name of style preset
        template_data: Optional template file data

    Returns:
        Configured PresentationView instance
    """
    view = PresentationView()

    if template_data:
        view.apply_template(template_data)

    # Select color scheme based on style
    color_schemes = {
        "professional": ColorScheme.PROFESSIONAL,
        "creative": ColorScheme.CREATIVE,
        "minimalist": ColorScheme.MINIMALIST,
        "technical": ColorScheme.TECHNICAL,
    }

    style = SlideStyle(
        color_scheme=color_schemes.get(style_name, ColorScheme.PROFESSIONAL)
    )
    view.set_style(style)

    return view
