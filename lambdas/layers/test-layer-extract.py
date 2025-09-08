#!/usr/bin/env python3
"""
Lambda Layer Extraction and Testing Script
Extracts the layer and tests imports in a simulated Lambda environment
"""

import sys
import os
import tempfile
import zipfile
import shutil
from pathlib import Path


def extract_and_test_layer():
    """Extract layer and test imports with proper Python path setup."""
    print("=== Lambda Layer Extraction and Testing ===")
    
    # Check if layer exists
    layer_path = Path("dist/ai-ppt-assistant-dependencies.zip")
    if not layer_path.exists():
        print("❌ Layer package not found. Run build script first.")
        return False
    
    print(f"✅ Found layer package: {layer_path}")
    print(f"Layer size: {layer_path.stat().st_size / (1024*1024):.2f} MB")
    
    # Create temporary directory for extraction
    with tempfile.TemporaryDirectory() as temp_dir:
        print(f"📦 Extracting layer to: {temp_dir}")
        
        # Extract the layer
        with zipfile.ZipFile(layer_path, 'r') as zip_ref:
            zip_ref.extractall(temp_dir)
        
        # Find the site-packages directory
        site_packages_path = None
        for root, dirs, files in os.walk(temp_dir):
            if root.endswith('site-packages'):
                site_packages_path = root
                break
        
        if not site_packages_path:
            print("❌ Could not find site-packages directory in layer")
            return False
        
        print(f"📂 Found site-packages: {site_packages_path}")
        
        # Add to Python path (at the beginning for priority)
        sys.path.insert(0, site_packages_path)
        
        # Test imports
        test_packages = [
            ('boto3', 'boto3'),
            ('aws-lambda-powertools', 'aws_lambda_powertools'),
            ('requests', 'requests'), 
            ('python-dateutil', 'dateutil'),
            ('python-dotenv', 'dotenv'),
            ('typing-extensions', 'typing_extensions'),
            ('cachetools', 'cachetools'),
            ('python-pptx', 'pptx'),
            ('Pillow', 'PIL'),
            ('jsonschema', 'jsonschema'),
            ('PyYAML', 'yaml'),
        ]
        
        success_count = 0
        total_count = len(test_packages)
        
        print("\n📋 Testing package imports:")
        print("-" * 50)
        
        for package_name, import_name in test_packages:
            try:
                __import__(import_name)
                print(f"✅ {package_name:20} → {import_name}")
                success_count += 1
            except ImportError as e:
                print(f"❌ {package_name:20} → {import_name} : {str(e)}")
            except Exception as e:
                print(f"⚠️  {package_name:20} → {import_name} : {str(e)}")
        
        print("-" * 50)
        print(f"Import Success Rate: {success_count}/{total_count} ({(success_count/total_count)*100:.1f}%)")
        
        # Test functionality
        print("\n🧪 Testing core functionality:")
        print("-" * 50)
        
        # Test AWS Lambda Powertools
        try:
            from aws_lambda_powertools import Logger, Tracer, Metrics
            logger = Logger()
            tracer = Tracer()  
            metrics = Metrics()
            print("✅ AWS Lambda Powertools: Logger, Tracer, Metrics initialized")
        except Exception as e:
            print(f"❌ AWS Lambda Powertools test failed: {e}")
        
        # Test PowerPoint creation
        try:
            from pptx import Presentation
            prs = Presentation()
            slide = prs.slides.add_slide(prs.slide_layouts[0])
            title = slide.shapes.title
            title.text = "Test Slide"
            print("✅ PowerPoint: Presentation creation successful")
        except Exception as e:
            print(f"❌ PowerPoint test failed: {e}")
        
        # Test image processing
        try:
            from PIL import Image
            import io
            img = Image.new('RGB', (100, 100), color='red')
            buffer = io.BytesIO()
            img.save(buffer, format='PNG')
            print("✅ Pillow: Image creation and saving successful")
        except Exception as e:
            print(f"❌ Pillow test failed: {e}")
        
        # Test JSON Schema validation
        try:
            import jsonschema
            schema = {"type": "object", "properties": {"name": {"type": "string"}}}
            jsonschema.validate({"name": "test"}, schema)
            print("✅ JSON Schema: Validation successful")
        except Exception as e:
            print(f"❌ JSON Schema test failed: {e}")
        
        # Test caching
        try:
            import cachetools
            cache = cachetools.LRUCache(maxsize=100)
            cache['test'] = 'value'
            assert cache['test'] == 'value'
            print("✅ Cachetools: LRU Cache test successful")
        except Exception as e:
            print(f"❌ Cachetools test failed: {e}")
        
        # Test YAML processing
        try:
            import yaml
            test_data = {'test': 'value', 'number': 42}
            yaml_str = yaml.dump(test_data)
            parsed = yaml.safe_load(yaml_str)
            assert parsed == test_data
            print("✅ PyYAML: YAML serialization/deserialization successful")
        except Exception as e:
            print(f"❌ PyYAML test failed: {e}")
        
        print("\n🔍 Layer Analysis:")
        print("-" * 50)
        
        # Count packages in site-packages
        package_dirs = [d for d in os.listdir(site_packages_path) 
                       if os.path.isdir(os.path.join(site_packages_path, d)) 
                       and not d.startswith('_')]
        print(f"📦 Total packages in layer: {len(package_dirs)}")
        
        # Show top-level packages
        print("🎯 Key packages found:")
        key_packages = ['boto3', 'botocore', 'aws_lambda_powertools', 'pptx', 'PIL', 'jsonschema', 'yaml', 'cachetools']
        found_key_packages = []
        
        for pkg in key_packages:
            pkg_path = os.path.join(site_packages_path, pkg)
            if os.path.exists(pkg_path):
                found_key_packages.append(pkg)
                print(f"   ✅ {pkg}")
            else:
                print(f"   ❌ {pkg}")
        
        print(f"\n📊 Summary:")
        print(f"   Layer size: {layer_path.stat().st_size / (1024*1024):.2f} MB")
        print(f"   Import success: {success_count}/{total_count} packages")
        print(f"   Key packages: {len(found_key_packages)}/{len(key_packages)}")
        
        # Overall assessment
        if success_count == total_count and len(found_key_packages) == len(key_packages):
            print("\n🎉 Layer is fully functional and ready for production!")
            return True
        elif success_count >= total_count * 0.8:
            print("\n⚠️  Layer is mostly functional but has some issues to address")
            return True
        else:
            print("\n❌ Layer has significant issues and needs to be rebuilt")
            return False


if __name__ == "__main__":
    print(f"Python version: {sys.version}")
    print(f"Working directory: {os.getcwd()}")
    print()
    
    success = extract_and_test_layer()
    sys.exit(0 if success else 1)