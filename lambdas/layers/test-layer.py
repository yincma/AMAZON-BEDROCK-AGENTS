#!/usr/bin/env python3
"""
Lambda Layer Compatibility Test Script
Tests if the built layer works correctly with AWS Lambda Python 3.12 runtime
"""

import sys
import os
import subprocess
import json
from pathlib import Path
from typing import Dict, List, Tuple, Optional
from datetime import datetime


def print_header(title: str) -> None:
    """Print a formatted header."""
    print(f"\n{'='*60}")
    print(f" {title}")
    print(f"{'='*60}")


def print_success(message: str) -> None:
    """Print success message in green."""
    print(f"✅ {message}")


def print_warning(message: str) -> None:
    """Print warning message in yellow."""
    print(f"⚠️  {message}")


def print_error(message: str) -> None:
    """Print error message in red."""
    print(f"❌ {message}")


def test_python_version() -> bool:
    """Test if we're running the correct Python version."""
    print_header("Python Version Check")
    
    version_info = sys.version_info
    python_version = f"{version_info.major}.{version_info.minor}"
    
    print(f"Current Python version: {sys.version}")
    
    if python_version == "3.12":
        print_success(f"Python version {python_version} matches Lambda runtime")
        return True
    else:
        print_warning(f"Python version {python_version} differs from Lambda runtime 3.12")
        print("This may cause compatibility issues in production")
        return False


def test_requirements_parsing() -> Tuple[bool, List[str]]:
    """Parse and validate requirements.txt."""
    print_header("Requirements Analysis")
    
    requirements_file = Path("requirements.txt")
    if not requirements_file.exists():
        print_error("requirements.txt not found")
        return False, []
    
    packages = []
    problematic_packages = []
    
    with open(requirements_file, 'r') as f:
        for line_num, line in enumerate(f, 1):
            line = line.strip()
            if line and not line.startswith('#'):
                if '==' in line:
                    package_name, version = line.split('==')
                    packages.append((package_name.strip(), version.strip()))
                    
                    # Check for known problematic versions
                    if package_name.strip() == "aws-lambda-powertools" and version.strip() == "2.39.0":
                        problematic_packages.append(f"Line {line_num}: {line} (known pydantic regression)")
                else:
                    print_warning(f"Line {line_num}: Unpinned version '{line}' may cause instability")
    
    print(f"Found {len(packages)} packages with pinned versions")
    
    if problematic_packages:
        print_warning("Potentially problematic packages found:")
        for pkg in problematic_packages:
            print(f"  - {pkg}")
    else:
        print_success("No known problematic package versions detected")
    
    # Display key packages
    key_packages = ["boto3", "aws-lambda-powertools", "python-pptx", "Pillow"]
    print("\nKey packages for this project:")
    for pkg_name, version in packages:
        if pkg_name in key_packages:
            print(f"  - {pkg_name}: {version}")
    
    return len(problematic_packages) == 0, [f"{name}=={version}" for name, version in packages]


def test_layer_structure() -> bool:
    """Test if the layer structure is correct."""
    print_header("Layer Structure Verification")
    
    dist_dir = Path("dist")
    if not dist_dir.exists():
        print_error("dist directory not found. Run build script first.")
        return False
    
    layer_zip = dist_dir / "ai-ppt-assistant-dependencies.zip"
    if not layer_zip.exists():
        print_error(f"Layer zip file not found: {layer_zip}")
        return False
    
    print_success(f"Layer zip file found: {layer_zip}")
    
    # Check file size
    size_bytes = layer_zip.stat().st_size
    size_mb = size_bytes / (1024 * 1024)
    
    print(f"Layer size: {size_mb:.2f} MB ({size_bytes:,} bytes)")
    
    if size_bytes > 50 * 1024 * 1024:  # 50MB
        print_warning("Layer size is quite large and may approach Lambda limits")
    else:
        print_success("Layer size is within reasonable limits")
    
    # Test zip integrity
    try:
        import zipfile
        with zipfile.ZipFile(layer_zip, 'r') as zf:
            bad_files = zf.testzip()
            if bad_files:
                print_error(f"Corrupted files in zip: {bad_files}")
                return False
            else:
                print_success("Zip file integrity verified")
            
            # Check structure
            file_list = zf.namelist()
            expected_structure = "python/lib/python3.12/site-packages/"
            
            if any(f.startswith(expected_structure) for f in file_list):
                print_success(f"Correct layer structure found: {expected_structure}")
            else:
                print_error(f"Expected structure '{expected_structure}' not found")
                print("Available top-level directories:")
                top_dirs = set(f.split('/')[0] for f in file_list if '/' in f)
                for d in sorted(top_dirs)[:10]:
                    print(f"  - {d}/")
                return False
                
    except Exception as e:
        print_error(f"Failed to test zip file: {e}")
        return False
    
    return True


def test_package_imports(packages: List[str]) -> Dict[str, bool]:
    """Test if packages can be imported successfully."""
    print_header("Package Import Testing")
    
    # Map package names to import names
    import_map = {
        "python-pptx": "pptx",
        "Pillow": "PIL",
        "aws-lambda-powertools": "aws_lambda_powertools",
        "python-dateutil": "dateutil",
        "PyYAML": "yaml",
    }
    
    results = {}
    critical_packages = ["boto3", "aws_lambda_powertools", "pptx", "PIL"]
    
    for package in packages:
        package_name = package.split('==')[0]
        import_name = import_map.get(package_name, package_name.replace('-', '_'))
        
        try:
            __import__(import_name)
            print_success(f"{package_name} → {import_name}")
            results[package_name] = True
        except ImportError as e:
            if package_name.replace('-', '_') in critical_packages or import_name in critical_packages:
                print_error(f"{package_name} → {import_name} (CRITICAL): {e}")
            else:
                print_warning(f"{package_name} → {import_name}: {e}")
            results[package_name] = False
        except Exception as e:
            print_error(f"{package_name} → {import_name} (ERROR): {e}")
            results[package_name] = False
    
    # Test specific functionality
    print("\nTesting specific functionality:")
    
    # Test AWS Lambda Powertools
    try:
        from aws_lambda_powertools import Logger, Tracer, Metrics
        logger = Logger()
        tracer = Tracer()
        metrics = Metrics()
        print_success("AWS Lambda Powertools components initialized")
    except Exception as e:
        print_error(f"AWS Lambda Powertools functionality test failed: {e}")
    
    # Test PowerPoint creation
    try:
        from pptx import Presentation
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        print_success("PowerPoint presentation creation test passed")
    except Exception as e:
        print_error(f"PowerPoint functionality test failed: {e}")
    
    # Test image processing
    try:
        from PIL import Image
        import io
        img = Image.new('RGB', (100, 100), color='red')
        buffer = io.BytesIO()
        img.save(buffer, format='PNG')
        print_success("Image processing test passed")
    except Exception as e:
        print_error(f"Image processing test failed: {e}")
    
    return results


def test_lambda_runtime_compatibility() -> bool:
    """Test Lambda runtime specific features."""
    print_header("Lambda Runtime Compatibility")
    
    # Test JSON serialization (common in Lambda)
    try:
        import json
        test_data = {"test": "value", "number": 42}
        serialized = json.dumps(test_data)
        deserialized = json.loads(serialized)
        assert deserialized == test_data
        print_success("JSON serialization works correctly")
    except Exception as e:
        print_error(f"JSON serialization test failed: {e}")
        return False
    
    # Test environment variable access
    try:
        import os
        env_test = os.environ.get("PATH", "not_found")
        if env_test != "not_found":
            print_success("Environment variable access works")
        else:
            print_warning("PATH environment variable not found (unusual)")
    except Exception as e:
        print_error(f"Environment variable test failed: {e}")
    
    # Test datetime handling (common in Lambda events)
    try:
        from datetime import datetime, timezone
        import json
        
        now = datetime.now(timezone.utc)
        iso_string = now.isoformat()
        parsed = datetime.fromisoformat(iso_string.replace('Z', '+00:00'))
        print_success("Datetime handling works correctly")
    except Exception as e:
        print_error(f"Datetime test failed: {e}")
    
    return True


def generate_report(test_results: Dict[str, any]) -> None:
    """Generate a test report."""
    print_header("Test Report Summary")
    
    total_tests = len(test_results)
    passed_tests = sum(1 for result in test_results.values() if result is True)
    
    print(f"Total tests: {total_tests}")
    print(f"Passed: {passed_tests}")
    print(f"Failed: {total_tests - passed_tests}")
    print(f"Success rate: {(passed_tests/total_tests)*100:.1f}%")
    
    if passed_tests == total_tests:
        print_success("All tests passed! Layer is ready for production use.")
    elif passed_tests >= total_tests * 0.8:
        print_warning("Most tests passed. Review warnings before production use.")
    else:
        print_error("Multiple test failures. Fix issues before production deployment.")
    
    # Save detailed report
    report_file = Path("dist/test-report.json")
    report_file.parent.mkdir(exist_ok=True)
    
    with open(report_file, 'w') as f:
        json.dump({
            "timestamp": datetime.now().isoformat(),
            "python_version": sys.version,
            "test_results": test_results,
            "summary": {
                "total_tests": total_tests,
                "passed_tests": passed_tests,
                "success_rate": (passed_tests/total_tests)*100
            }
        }, f, indent=2)
    
    print(f"\nDetailed report saved to: {report_file}")


def main():
    """Run all tests."""
    print_header("Lambda Layer Compatibility Test Suite")
    print(f"Testing environment: Python {sys.version}")
    print(f"Working directory: {os.getcwd()}")
    
    test_results = {}
    
    # Run tests
    test_results["python_version"] = test_python_version()
    
    requirements_ok, packages = test_requirements_parsing()
    test_results["requirements_parsing"] = requirements_ok
    
    test_results["layer_structure"] = test_layer_structure()
    
    if packages:
        import_results = test_package_imports(packages)
        test_results["package_imports"] = all(import_results.values())
        test_results["individual_imports"] = import_results
    else:
        test_results["package_imports"] = False
        test_results["individual_imports"] = {}
    
    test_results["lambda_compatibility"] = test_lambda_runtime_compatibility()
    
    # Generate report
    generate_report(test_results)
    
    # Exit with appropriate code
    if all(result is True for key, result in test_results.items() 
           if key not in ["individual_imports"]):
        sys.exit(0)
    else:
        sys.exit(1)


if __name__ == "__main__":
    main()