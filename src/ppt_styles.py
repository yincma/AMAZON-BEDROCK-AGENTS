"""
PPT样式管理 - 定义演示文稿的视觉样式和格式
"""

from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor


class PPTStyles:
    """PPT样式配置类"""

    # 字体配置
    TITLE_FONT_SIZE = Pt(44)
    SUBTITLE_FONT_SIZE = Pt(32)
    HEADING_FONT_SIZE = Pt(36)
    BODY_FONT_SIZE = Pt(24)

    # 颜色配置
    PRIMARY_COLOR = RGBColor(0, 51, 102)      # 深蓝色
    SECONDARY_COLOR = RGBColor(51, 51, 51)    # 深灰色
    ACCENT_COLOR = RGBColor(255, 102, 0)      # 橙色
    WHITE_COLOR = RGBColor(255, 255, 255)     # 白色
    LIGHT_GRAY = RGBColor(240, 240, 240)      # 浅灰色

    # 边距配置
    MARGIN_TOP = Inches(1)
    MARGIN_LEFT = Inches(1)
    MARGIN_RIGHT = Inches(1)
    MARGIN_BOTTOM = Inches(1)

    # 间距配置
    PARAGRAPH_SPACING = Pt(6)
    LINE_SPACING = 1.2

    @staticmethod
    def apply_title_style(shape):
        """
        应用标题样式

        Args:
            shape: 形状对象，应包含文本框

        Returns:
            paragraph: 配置好的段落对象
        """
        if shape.has_text_frame:
            text_frame = shape.text_frame
            text_frame.clear()

            # 创建段落
            p = text_frame.add_paragraph()
            p.alignment = PP_ALIGN.CENTER

            # 设置字体样式
            run = p.runs[0] if p.runs else None
            if run:
                font = run.font
                font.size = PPTStyles.TITLE_FONT_SIZE
                font.bold = True
                font.color.rgb = PPTStyles.PRIMARY_COLOR

            return p

        return None

    @staticmethod
    def apply_subtitle_style(shape):
        """
        应用副标题样式

        Args:
            shape: 形状对象

        Returns:
            paragraph: 配置好的段落对象
        """
        if shape.has_text_frame:
            text_frame = shape.text_frame
            text_frame.clear()

            p = text_frame.add_paragraph()
            p.alignment = PP_ALIGN.CENTER

            # 设置字体样式
            run = p.runs[0] if p.runs else None
            if run:
                font = run.font
                font.size = PPTStyles.SUBTITLE_FONT_SIZE
                font.color.rgb = PPTStyles.SECONDARY_COLOR

            return p

        return None

    @staticmethod
    def apply_bullet_style(text_frame):
        """
        应用要点样式

        Args:
            text_frame: 文本框对象
        """
        # 设置文本框边距
        text_frame.margin_left = Inches(0.5)
        text_frame.margin_top = Inches(0.3)
        text_frame.margin_right = Inches(0.3)
        text_frame.margin_bottom = Inches(0.3)

        # 应用段落样式
        for paragraph in text_frame.paragraphs:
            paragraph.level = 0  # 第一级要点
            paragraph.space_after = PPTStyles.PARAGRAPH_SPACING

            # 设置字体
            for run in paragraph.runs:
                font = run.font
                font.size = PPTStyles.BODY_FONT_SIZE
                font.color.rgb = PPTStyles.SECONDARY_COLOR

    @staticmethod
    def apply_heading_style(shape):
        """
        应用标题样式（内容页）

        Args:
            shape: 标题形状对象

        Returns:
            paragraph: 配置好的段落对象
        """
        if shape.has_text_frame:
            text_frame = shape.text_frame

            # 获取第一个段落或创建新段落
            if text_frame.paragraphs:
                p = text_frame.paragraphs[0]
            else:
                p = text_frame.add_paragraph()

            p.alignment = PP_ALIGN.LEFT

            # 设置字体样式
            for run in p.runs:
                font = run.font
                font.size = PPTStyles.HEADING_FONT_SIZE
                font.bold = True
                font.color.rgb = PPTStyles.PRIMARY_COLOR

            return p

        return None

    @staticmethod
    def apply_speaker_notes_style(notes_slide):
        """
        应用演讲者备注样式

        Args:
            notes_slide: 备注页对象
        """
        for shape in notes_slide.shapes:
            if shape.has_text_frame:
                text_frame = shape.text_frame

                for paragraph in text_frame.paragraphs:
                    for run in paragraph.runs:
                        font = run.font
                        font.size = Pt(14)  # 备注文字较小
                        font.color.rgb = PPTStyles.SECONDARY_COLOR

    @staticmethod
    def configure_slide_layout(slide):
        """
        配置幻灯片布局

        Args:
            slide: 幻灯片对象
        """
        # 可以在这里设置背景色、边框等
        pass

    @staticmethod
    def get_color_scheme():
        """
        获取完整的颜色方案

        Returns:
            dict: 颜色配置字典
        """
        return {
            'primary': PPTStyles.PRIMARY_COLOR,
            'secondary': PPTStyles.SECONDARY_COLOR,
            'accent': PPTStyles.ACCENT_COLOR,
            'white': PPTStyles.WHITE_COLOR,
            'light_gray': PPTStyles.LIGHT_GRAY
        }

    @staticmethod
    def get_font_sizes():
        """
        获取字体大小配置

        Returns:
            dict: 字体大小配置字典
        """
        return {
            'title': PPTStyles.TITLE_FONT_SIZE,
            'subtitle': PPTStyles.SUBTITLE_FONT_SIZE,
            'heading': PPTStyles.HEADING_FONT_SIZE,
            'body': PPTStyles.BODY_FONT_SIZE
        }


class TemplateManager:
    """模板管理器"""

    AVAILABLE_TEMPLATES = ['default', 'corporate', 'academic', 'creative']

    @staticmethod
    def apply_template(presentation, template_name: str = 'default'):
        """
        应用指定模板

        Args:
            presentation: 演示文稿对象
            template_name: 模板名称

        Raises:
            ValueError: 当模板不存在时
        """
        if template_name not in TemplateManager.AVAILABLE_TEMPLATES:
            raise ValueError(f"Template '{template_name}' not available. "
                           f"Available templates: {TemplateManager.AVAILABLE_TEMPLATES}")

        if template_name == 'default':
            TemplateManager._apply_default_template(presentation)
        elif template_name == 'corporate':
            TemplateManager._apply_corporate_template(presentation)
        elif template_name == 'academic':
            TemplateManager._apply_academic_template(presentation)
        elif template_name == 'creative':
            TemplateManager._apply_creative_template(presentation)

    @staticmethod
    def _apply_default_template(presentation):
        """应用默认模板"""
        # 默认模板使用标准样式，无需特殊处理
        pass

    @staticmethod
    def _apply_corporate_template(presentation):
        """应用企业模板"""
        # 企业模板：使用深蓝色主题，正式字体
        PPTStyles.PRIMARY_COLOR = RGBColor(0, 51, 102)
        PPTStyles.SECONDARY_COLOR = RGBColor(51, 51, 51)

    @staticmethod
    def _apply_academic_template(presentation):
        """应用学术模板"""
        # 学术模板：使用深绿色主题，serif字体
        PPTStyles.PRIMARY_COLOR = RGBColor(0, 102, 51)
        PPTStyles.SECONDARY_COLOR = RGBColor(51, 51, 51)

    @staticmethod
    def _apply_creative_template(presentation):
        """应用创意模板"""
        # 创意模板：使用紫色主题，现代字体
        PPTStyles.PRIMARY_COLOR = RGBColor(102, 51, 153)
        PPTStyles.ACCENT_COLOR = RGBColor(255, 153, 0)