"""
PPT编译器 - 将JSON内容编译为PPTX文件
使用python-pptx库生成演示文稿
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
import boto3
import json
import os
import tempfile
import logging
from typing import Dict, List
from io import BytesIO
import requests

logger = logging.getLogger(__name__)
logger.setLevel(logging.INFO)


def create_pptx_from_content(content: Dict, include_notes: bool = True) -> bytes:
    """
    从JSON内容创建PPTX文件

    Args:
        content: 包含slides列表的字典
        include_notes: 是否包含演讲者备注

    Returns:
        bytes: PPTX文件的字节数据
    """
    if not content or 'slides' not in content:
        raise ValueError("Content must contain 'slides' key")

    slides = content['slides']
    if not slides:
        raise ValueError("Content must contain at least one slide")

    # 创建新的演示文稿
    prs = Presentation()

    # 为每个幻灯片添加内容
    for slide_data in slides:
        add_content_slide(prs, slide_data, include_notes)

    # 将演示文稿保存到字节流
    pptx_bytes = BytesIO()
    prs.save(pptx_bytes)
    pptx_bytes.seek(0)

    return pptx_bytes.getvalue()


def add_content_slide(prs: Presentation, slide_data: Dict, include_notes: bool = True):
    """
    添加内容页到演示文稿

    Args:
        prs: 演示文稿对象
        slide_data: 幻灯片数据，包含title, bullet_points和image_url
        include_notes: 是否包含演讲者备注
    """
    # 检查是否有图片URL
    has_image = 'image_url' in slide_data and slide_data['image_url']

    if has_image:
        # 使用空白布局以便自定义放置元素
        slide_layout = prs.slide_layouts[5]  # 布局5通常是空白布局
        slide = prs.slides.add_slide(slide_layout)

        # 手动添加标题文本框
        left = Inches(0.5)
        top = Inches(0.5)
        width = Inches(9)
        height = Inches(1)
        title_box = slide.shapes.add_textbox(left, top, width, height)
        title_frame = title_box.text_frame
        if 'title' in slide_data:
            title_p = title_frame.paragraphs[0]
            title_p.text = slide_data['title']
            title_p.font.size = Pt(32)
            title_p.font.bold = True

        # 左侧添加要点内容
        if 'bullet_points' in slide_data:
            left = Inches(0.5)
            top = Inches(1.8)
            width = Inches(4.5)
            height = Inches(4.5)
            content_box = slide.shapes.add_textbox(left, top, width, height)
            text_frame = content_box.text_frame

            for i, bullet_point in enumerate(slide_data['bullet_points']):
                if i == 0:
                    p = text_frame.paragraphs[0]
                else:
                    p = text_frame.add_paragraph()
                p.text = f"• {bullet_point}"
                p.font.size = Pt(16)
                p.space_after = Pt(12)

        # 右侧添加图片
        try:
            # 处理图片URL
            image_url = slide_data['image_url']
            if image_url.startswith('http'):
                # 下载HTTP/HTTPS图片
                response = requests.get(image_url, timeout=10)
                if response.status_code == 200:
                    image_stream = BytesIO(response.content)
                    left = Inches(5.5)
                    top = Inches(1.8)
                    width = Inches(4)
                    height = Inches(3)
                    slide.shapes.add_picture(image_stream, left, top, width, height)
                    logger.info(f"成功添加图片到幻灯片")
            else:
                # 假设是S3 URL或本地路径
                # 从S3获取图片
                s3 = boto3.client('s3')
                # 从完整URL解析bucket和key
                if 's3.amazonaws.com' in image_url:
                    # 移除协议前缀
                    url_without_protocol = image_url.replace('https://', '').replace('http://', '')
                    parts = url_without_protocol.split('/')
                    domain_part = parts[0]

                    # 处理不同的S3 URL格式
                    if '.s3.' in domain_part:
                        # 格式: bucket-name.s3.region.amazonaws.com 或 bucket-name.s3.amazonaws.com
                        bucket_with_suffix = domain_part.split('.s3.')[0]

                        # 移除AWS账号ID后缀（如果存在）
                        # AWS账号ID通常是12位数字，但为了兼容性检查10位以上的数字
                        if '-' in bucket_with_suffix:
                            parts_bucket = bucket_with_suffix.split('-')
                            if parts_bucket[-1].isdigit() and len(parts_bucket[-1]) >= 10:
                                bucket_name = '-'.join(parts_bucket[:-1])
                            else:
                                bucket_name = bucket_with_suffix
                        else:
                            bucket_name = bucket_with_suffix
                    else:
                        # 回退到原始逻辑
                        bucket_name = domain_part.split('.')[0]

                    key = '/'.join(parts[1:])
                else:
                    # 可能是相对路径，跳过图片
                    logger.warning(f"无法解析图片URL: {image_url}")
                    raise ValueError("无法解析图片URL")

                obj = s3.get_object(Bucket=bucket_name, Key=key)
                image_stream = BytesIO(obj['Body'].read())
                left = Inches(5.5)
                top = Inches(1.8)
                width = Inches(4)
                height = Inches(3)
                slide.shapes.add_picture(image_stream, left, top, width, height)
                logger.info(f"成功从S3添加图片到幻灯片")
        except Exception as e:
            logger.warning(f"添加图片失败: {str(e)}")
            # 即使图片添加失败，也继续处理
    else:
        # 没有图片，使用标准布局
        slide_layout = prs.slide_layouts[1]  # 布局1通常是标题+内容
        slide = prs.slides.add_slide(slide_layout)

        # 添加标题
        title = slide.shapes.title
        if 'title' in slide_data:
            title.text = slide_data['title']

        # 添加要点内容
        if 'bullet_points' in slide_data:
            content = slide.placeholders[1]  # 内容占位符
            text_frame = content.text_frame
            text_frame.clear()  # 清除默认文本

            # 添加每个要点
            for i, bullet_point in enumerate(slide_data['bullet_points']):
                if i == 0:
                    # 第一个段落已存在，直接使用
                    p = text_frame.paragraphs[0]
                else:
                    # 添加新段落
                    p = text_frame.add_paragraph()

                p.text = bullet_point
                p.level = 0  # 设置为第一级要点

    # 添加演讲者备注
    if include_notes and 'speaker_notes' in slide_data:
        notes_slide = slide.notes_slide
        text_frame = notes_slide.notes_text_frame
        text_frame.text = slide_data['speaker_notes']


def get_slide_count(pptx_bytes: bytes) -> int:
    """
    获取PPTX文件的幻灯片数量

    Args:
        pptx_bytes: PPTX文件字节数据

    Returns:
        int: 幻灯片数量
    """
    pptx_stream = BytesIO(pptx_bytes)
    prs = Presentation(pptx_stream)
    return len(prs.slides)


def extract_text_content(pptx_bytes: bytes) -> str:
    """
    从PPTX文件中提取文本内容

    Args:
        pptx_bytes: PPTX文件字节数据

    Returns:
        str: 提取的文本内容
    """
    pptx_stream = BytesIO(pptx_bytes)
    prs = Presentation(pptx_stream)

    text_content = []

    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text_content.append(shape.text)

    return " ".join(text_content)


def create_pptx_with_template(content: Dict, template_name: str = "default") -> bytes:
    """
    使用指定模板创建PPTX文件

    Args:
        content: 幻灯片内容
        template_name: 模板名称

    Returns:
        bytes: PPTX文件字节数据
    """
    if template_name != "default":
        raise FileNotFoundError(f"Template '{template_name}' not found")

    # 目前只支持默认模板
    return create_pptx_from_content(content)


def extract_speaker_notes(pptx_bytes: bytes) -> List[str]:
    """
    提取演讲者备注

    Args:
        pptx_bytes: PPTX文件字节数据

    Returns:
        List[str]: 每页的演讲者备注列表
    """
    pptx_stream = BytesIO(pptx_bytes)
    prs = Presentation(pptx_stream)

    notes = []
    for slide in prs.slides:
        if slide.has_notes_slide:
            notes_slide = slide.notes_slide
            # 提取备注文本
            notes_text = ""
            for shape in notes_slide.shapes:
                if hasattr(shape, "text"):
                    notes_text += shape.text
            notes.append(notes_text)
        else:
            notes.append("")

    return notes


class PPTCompiler:
    """PPT编译器类 - 处理S3集成和文件管理"""

    def __init__(self):
        self.s3_client = boto3.client('s3')
        self.bucket_name = os.environ.get('S3_BUCKET', 'ai-ppt-presentations-dev')

    def compile_ppt(self, presentation_id: str) -> str:
        """
        编译PPT的主函数

        Args:
            presentation_id: 演示文稿ID

        Returns:
            str: 预签名下载URL
        """
        # 从S3读取内容
        content_key = f"presentations/{presentation_id}/content.json"

        try:
            response = self.s3_client.get_object(
                Bucket=self.bucket_name,
                Key=content_key
            )
            content = json.loads(response['Body'].read().decode('utf-8'))
        except Exception as e:
            logger.error(f"Failed to read content from S3: {e}")
            raise

        # 生成PPTX
        pptx_bytes = create_pptx_from_content(content)

        # 保存到S3
        s3_key = self.save_to_s3(pptx_bytes, presentation_id)

        # 生成下载链接
        download_url = self.generate_download_url(s3_key)

        return download_url

    def save_to_s3(self, pptx_bytes: bytes, presentation_id: str) -> str:
        """
        保存PPTX到S3

        Args:
            pptx_bytes: PPTX文件字节数据
            presentation_id: 演示文稿ID

        Returns:
            str: S3对象键
        """
        s3_key = f"presentations/{presentation_id}/output/presentation.pptx"

        try:
            self.s3_client.put_object(
                Bucket=self.bucket_name,
                Key=s3_key,
                Body=pptx_bytes,
                ContentType='application/vnd.openxmlformats-officedocument.presentationml.presentation'
            )
            logger.info(f"Successfully saved PPTX to S3: {s3_key}")
            return s3_key
        except Exception as e:
            logger.error(f"Failed to save PPTX to S3: {e}")
            raise

    def generate_download_url(self, s3_key: str, expires_in: int = 3600) -> str:
        """
        生成预签名下载URL

        Args:
            s3_key: S3对象键
            expires_in: URL有效期（秒）

        Returns:
            str: 预签名URL
        """
        try:
            url = self.s3_client.generate_presigned_url(
                'get_object',
                Params={'Bucket': self.bucket_name, 'Key': s3_key},
                ExpiresIn=expires_in
            )
            return url
        except Exception as e:
            logger.error(f"Failed to generate download URL: {e}")
            raise


# S3集成函数 - 与测试兼容
def save_pptx_to_s3(pptx_bytes: bytes, presentation_id: str, s3_client, bucket_name: str) -> str:
    """
    保存PPTX到S3（兼容测试接口）

    Args:
        pptx_bytes: PPTX字节数据
        presentation_id: 演示文稿ID
        s3_client: S3客户端
        bucket_name: 存储桶名称

    Returns:
        str: S3对象键
    """
    s3_key = f"presentations/{presentation_id}/output/presentation.pptx"

    s3_client.put_object(
        Bucket=bucket_name,
        Key=s3_key,
        Body=pptx_bytes,
        ContentType='application/vnd.openxmlformats-officedocument.presentationml.presentation'
    )

    return s3_key


def generate_download_url(presentation_id: str, s3_client, bucket_name: str, expires_in: int = 3600) -> str:
    """
    生成预签名下载URL（兼容测试接口）
    """
    s3_key = f"presentations/{presentation_id}/output/presentation.pptx"

    url = s3_client.generate_presigned_url(
        'get_object',
        Params={'Bucket': bucket_name, 'Key': s3_key},
        ExpiresIn=expires_in
    )

    return url


def save_pptx_with_metadata(pptx_bytes: bytes, metadata: Dict, presentation_id: str, s3_client, bucket_name: str):
    """
    保存PPTX文件和元数据到S3

    Args:
        pptx_bytes: PPTX字节数据
        metadata: 元数据字典
        presentation_id: 演示文稿ID
        s3_client: S3客户端
        bucket_name: 存储桶名称
    """
    # 保存PPTX文件
    pptx_key = f"presentations/{presentation_id}/output/presentation.pptx"
    s3_client.put_object(
        Bucket=bucket_name,
        Key=pptx_key,
        Body=pptx_bytes,
        ContentType='application/vnd.openxmlformats-officedocument.presentationml.presentation'
    )

    # 保存元数据
    metadata_key = f"presentations/{presentation_id}/metadata.json"
    s3_client.put_object(
        Bucket=bucket_name,
        Key=metadata_key,
        Body=json.dumps(metadata).encode('utf-8'),
        ContentType='application/json'
    )