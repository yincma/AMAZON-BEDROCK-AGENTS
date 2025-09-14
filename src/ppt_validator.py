"""
PPT验证器 - 验证生成的PPTX文件的完整性和质量
"""

import zipfile
import logging
from typing import Dict, Any, List
from io import BytesIO
from pptx import Presentation

logger = logging.getLogger(__name__)


def validate_pptx_integrity(pptx_bytes: bytes) -> bool:
    """
    验证PPTX文件完整性

    Args:
        pptx_bytes: PPTX文件字节数据

    Returns:
        bool: 文件是否完整有效
    """
    try:
        # 检查是否为有效的ZIP文件（PPTX本质上是ZIP）
        if not zipfile.is_zipfile(BytesIO(pptx_bytes)):
            logger.error("PPTX file is not a valid ZIP file")
            return False

        # 尝试使用python-pptx打开文件
        pptx_stream = BytesIO(pptx_bytes)
        prs = Presentation(pptx_stream)

        # 检查是否至少有一个幻灯片
        if len(prs.slides) == 0:
            logger.error("PPTX file contains no slides")
            return False

        logger.info(f"PPTX file validation passed: {len(prs.slides)} slides found")
        return True

    except Exception as e:
        logger.error(f"PPTX file validation failed: {e}")
        return False


def validate_content_completeness(input_content: Dict[str, Any], pptx_bytes: bytes) -> bool:
    """
    验证幻灯片内容完整性

    Args:
        input_content: 输入的内容数据
        pptx_bytes: 生成的PPTX文件字节数据

    Returns:
        bool: 内容是否完整
    """
    try:
        # 提取输入内容的文本
        input_texts = []
        if 'slides' in input_content:
            for slide in input_content['slides']:
                if 'title' in slide:
                    input_texts.append(slide['title'])
                if 'bullet_points' in slide:
                    input_texts.extend(slide['bullet_points'])

        # 从PPTX文件提取文本
        pptx_stream = BytesIO(pptx_bytes)
        prs = Presentation(pptx_stream)

        pptx_text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    pptx_text += shape.text + " "

        # 检查输入文本是否都存在于PPTX中
        missing_content = []
        for text in input_texts:
            if text.strip() and text.strip() not in pptx_text:
                missing_content.append(text)

        if missing_content:
            logger.warning(f"Missing content in PPTX: {missing_content}")
            return False

        logger.info("Content completeness validation passed")
        return True

    except Exception as e:
        logger.error(f"Content completeness validation failed: {e}")
        return False


def validate_accessibility(pptx_bytes: bytes) -> float:
    """
    验证PPTX可访问性

    Args:
        pptx_bytes: PPTX文件字节数据

    Returns:
        float: 可访问性得分 (0.0-1.0)
    """
    try:
        pptx_stream = BytesIO(pptx_bytes)
        prs = Presentation(pptx_stream)

        total_checks = 0
        passed_checks = 0

        for slide in prs.slides:
            # 检查1: 是否有标题
            total_checks += 1
            has_title = False

            for shape in slide.shapes:
                if shape.shape_type == 1 and hasattr(shape, 'text'):  # 文本框
                    if shape.text.strip():
                        has_title = True
                        break

            if has_title:
                passed_checks += 1

            # 检查2: 文本是否可读（非空）
            total_checks += 1
            has_readable_text = False

            for shape in slide.shapes:
                if hasattr(shape, 'text') and len(shape.text.strip()) > 5:
                    has_readable_text = True
                    break

            if has_readable_text:
                passed_checks += 1

        # 计算可访问性得分
        if total_checks == 0:
            return 0.0

        score = passed_checks / total_checks
        logger.info(f"Accessibility score: {score:.2f} ({passed_checks}/{total_checks})")

        return score

    except Exception as e:
        logger.error(f"Accessibility validation failed: {e}")
        return 0.0


def validate_file_size(pptx_bytes: bytes, max_size_mb: int = 100) -> bool:
    """
    验证文件大小限制

    Args:
        pptx_bytes: PPTX文件字节数据
        max_size_mb: 最大文件大小（MB）

    Returns:
        bool: 文件大小是否在限制内
    """
    try:
        size_mb = len(pptx_bytes) / (1024 * 1024)

        if size_mb > max_size_mb:
            logger.warning(f"PPTX file size too large: {size_mb:.2f} MB (max: {max_size_mb} MB)")
            return False

        logger.info(f"PPTX file size validation passed: {size_mb:.2f} MB")
        return True

    except Exception as e:
        logger.error(f"File size validation failed: {e}")
        return False


def validate_slide_structure(pptx_bytes: bytes) -> Dict[str, Any]:
    """
    验证幻灯片结构

    Args:
        pptx_bytes: PPTX文件字节数据

    Returns:
        Dict: 结构验证结果
    """
    try:
        pptx_stream = BytesIO(pptx_bytes)
        prs = Presentation(pptx_stream)

        validation_result = {
            'total_slides': len(prs.slides),
            'slides_with_title': 0,
            'slides_with_content': 0,
            'slides_with_images': 0,
            'empty_slides': 0,
            'valid': True,
            'issues': []
        }

        for i, slide in enumerate(prs.slides):
            slide_has_title = False
            slide_has_content = False
            slide_has_images = False
            text_shapes = 0

            for shape in slide.shapes:
                if hasattr(shape, 'text'):
                    if shape.text.strip():
                        text_shapes += 1
                        if not slide_has_title:
                            slide_has_title = True
                        elif not slide_has_content:
                            slide_has_content = True

                # 检查图片（简化检查）
                if shape.shape_type == 13:  # Picture
                    slide_has_images = True

            # 统计
            if slide_has_title:
                validation_result['slides_with_title'] += 1
            if slide_has_content:
                validation_result['slides_with_content'] += 1
            if slide_has_images:
                validation_result['slides_with_images'] += 1

            # 检查空幻灯片
            if text_shapes == 0:
                validation_result['empty_slides'] += 1
                validation_result['issues'].append(f"Slide {i+1}: No text content")

            # 检查只有标题没有内容的幻灯片
            if slide_has_title and not slide_has_content:
                validation_result['issues'].append(f"Slide {i+1}: Title only, no content")

        # 判断整体有效性
        if validation_result['empty_slides'] > 0:
            validation_result['valid'] = False

        logger.info(f"Slide structure validation completed: {validation_result}")
        return validation_result

    except Exception as e:
        logger.error(f"Slide structure validation failed: {e}")
        return {
            'valid': False,
            'error': str(e)
        }


class PPTXValidator:
    """PPTX验证器类"""

    def __init__(self):
        self.validation_rules = {
            'min_slides': 1,
            'max_slides': 50,
            'max_file_size_mb': 100,
            'min_accessibility_score': 0.7
        }

    def validate_all(self, pptx_bytes: bytes, input_content: Dict[str, Any] = None) -> Dict[str, Any]:
        """
        执行全面验证

        Args:
            pptx_bytes: PPTX文件字节数据
            input_content: 可选的输入内容，用于内容完整性检查

        Returns:
            Dict: 完整的验证结果
        """
        results = {
            'overall_valid': True,
            'checks': {},
            'issues': [],
            'warnings': []
        }

        # 1. 文件完整性检查
        integrity_valid = validate_pptx_integrity(pptx_bytes)
        results['checks']['integrity'] = integrity_valid
        if not integrity_valid:
            results['overall_valid'] = False
            results['issues'].append("File integrity check failed")

        # 2. 文件大小检查
        size_valid = validate_file_size(pptx_bytes, self.validation_rules['max_file_size_mb'])
        results['checks']['file_size'] = size_valid
        if not size_valid:
            results['overall_valid'] = False
            results['issues'].append("File size exceeds limit")

        # 3. 结构检查
        structure_result = validate_slide_structure(pptx_bytes)
        results['checks']['structure'] = structure_result

        if not structure_result.get('valid', False):
            results['overall_valid'] = False
            results['issues'].append("Slide structure validation failed")

        # 检查幻灯片数量
        slide_count = structure_result.get('total_slides', 0)
        if slide_count < self.validation_rules['min_slides']:
            results['overall_valid'] = False
            results['issues'].append(f"Too few slides: {slide_count}")
        elif slide_count > self.validation_rules['max_slides']:
            results['warnings'].append(f"Many slides: {slide_count}")

        # 4. 可访问性检查
        accessibility_score = validate_accessibility(pptx_bytes)
        results['checks']['accessibility'] = accessibility_score

        if accessibility_score < self.validation_rules['min_accessibility_score']:
            results['warnings'].append(f"Low accessibility score: {accessibility_score:.2f}")

        # 5. 内容完整性检查（如果提供了输入内容）
        if input_content:
            content_complete = validate_content_completeness(input_content, pptx_bytes)
            results['checks']['content_completeness'] = content_complete

            if not content_complete:
                results['overall_valid'] = False
                results['issues'].append("Content completeness check failed")

        logger.info(f"Full validation completed. Overall valid: {results['overall_valid']}")
        return results

    def update_validation_rules(self, new_rules: Dict[str, Any]):
        """
        更新验证规则

        Args:
            new_rules: 新的验证规则
        """
        self.validation_rules.update(new_rules)
        logger.info(f"Updated validation rules: {self.validation_rules}")