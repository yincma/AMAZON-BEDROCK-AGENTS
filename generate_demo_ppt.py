#!/usr/bin/env python3
"""
生成完整的AI PPT演示文稿
"""

import sys
import os
sys.path.append('lambdas')

from datetime import datetime
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from image_processing_service import ImageProcessingService
import io

def generate_demo_ppt():
    """生成包含AI图片的演示PPT"""

    print("🎨 开始生成AI驱动的PPT演示...")

    # 创建演示文稿
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # 图片生成服务
    image_service = ImageProcessingService()

    # 幻灯片内容定义
    slides_content = [
        {
            "title": "AI技术革新2025",
            "content": [
                "大语言模型的突破性进展",
                "多模态AI的商业化应用",
                "AGI时代的来临"
            ],
            "subtitle": "引领未来的智能技术"
        },
        {
            "title": "机器学习核心技术",
            "content": [
                "深度学习神经网络",
                "强化学习算法优化",
                "迁移学习与少样本学习"
            ],
            "subtitle": "构建智能系统的基石"
        },
        {
            "title": "AI应用场景",
            "content": [
                "智慧医疗与精准诊断",
                "智能制造与工业4.0",
                "金融科技与风险控制"
            ],
            "subtitle": "改变世界的实际应用"
        },
        {
            "title": "未来展望",
            "content": [
                "通用人工智能(AGI)发展",
                "人机协作新模式",
                "AI伦理与监管框架"
            ],
            "subtitle": "迈向智能新纪元"
        }
    ]

    # 为每张幻灯片生成内容和图片
    for i, slide_data in enumerate(slides_content, 1):
        print(f"\n📄 生成第 {i}/{len(slides_content)} 张幻灯片: {slide_data['title']}")

        # 添加幻灯片
        slide_layout = prs.slide_layouts[5]  # 空白布局
        slide = prs.slides.add_slide(slide_layout)

        # 添加标题
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.5),
            Inches(9), Inches(1)
        )
        title_frame = title_box.text_frame
        title_frame.text = slide_data["title"]
        title_frame.paragraphs[0].font.size = Pt(36)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        # 添加副标题
        subtitle_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(1.3),
            Inches(9), Inches(0.5)
        )
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.text = slide_data["subtitle"]
        subtitle_frame.paragraphs[0].font.size = Pt(18)
        subtitle_frame.paragraphs[0].font.italic = True
        subtitle_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        # 生成AI图片
        print(f"  🎨 生成AI图片...")
        prompt = image_service.generate_prompt(slide_data, "business")

        try:
            image_data = image_service.call_image_generation(prompt)

            if image_data and len(image_data) > 100000:  # 真实AI图片
                print(f"  ✅ AI图片生成成功 ({len(image_data):,} bytes)")
                # 添加图片到幻灯片
                image_stream = io.BytesIO(image_data)
                slide.shapes.add_picture(
                    image_stream,
                    Inches(5.5), Inches(2),
                    width=Inches(4), height=Inches(3)
                )
            else:
                print(f"  ⚠️ 使用占位图")
                # 使用占位图
                if image_data:
                    image_stream = io.BytesIO(image_data)
                    slide.shapes.add_picture(
                        image_stream,
                        Inches(5.5), Inches(2),
                        width=Inches(4), height=Inches(3)
                    )
        except Exception as e:
            print(f"  ❌ 图片生成失败: {e}")

        # 添加内容要点
        content_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(2),
            Inches(4.5), Inches(4)
        )
        content_frame = content_box.text_frame

        for point in slide_data["content"]:
            p = content_frame.add_paragraph()
            p.text = f"• {point}"
            p.font.size = Pt(16)
            p.level = 0
            p.space_after = Pt(12)

        # 添加页脚
        footer_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(6.8),
            Inches(9), Inches(0.5)
        )
        footer_frame = footer_box.text_frame
        footer_frame.text = f"AI-PPT Assistant | {datetime.now().strftime('%Y-%m-%d')} | 第 {i} 页"
        footer_frame.paragraphs[0].font.size = Pt(10)
        footer_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # 保存演示文稿
    filename = f"AI_Demo_PPT_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"
    prs.save(filename)

    print(f"\n✅ PPT生成完成！")
    print(f"📁 文件保存为: {filename}")
    print(f"📊 共 {len(slides_content)} 张幻灯片")

    return filename

if __name__ == "__main__":
    try:
        filename = generate_demo_ppt()
        print(f"\n🎉 成功生成演示文稿: {filename}")
        print("💡 提示: 可以用PowerPoint或Keynote打开查看")
    except Exception as e:
        print(f"\n❌ 生成失败: {e}")
        import traceback
        traceback.print_exc()